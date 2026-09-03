"""Build an editable non-oil report from the approved PowerPoint template.

The final PPTX is the source of truth. Browser preview images are rendered from
that PPTX so the preview and downloaded deck share one layout and pagination.
"""

from __future__ import annotations

import math
import os
import re
import shutil
import subprocess
import tempfile
import uuid
from copy import deepcopy
from pathlib import Path
from zipfile import ZipFile

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_DATA_LABEL_POSITION,
    XL_LEGEND_POSITION,
    XL_TICK_LABEL_POSITION,
)
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.opc.packuri import PackURI
from pptx.util import Inches, Pt


TEMPLATE_FILE = (
    Path(__file__).resolve().parent
    / "assets"
    / "non_oil_report_template"
    / "non_oil_report_template.pptx"
)
CANVAS_SIZE = (2560, 1440)
UNIT_ORDER = [
    "浦东", "闵普徐", "松金", "嘉青", "南汇", "宝静", "奉贤", "崇明",
    "中油奉贤", "中油同盛", "中油康桥", "中油农工商", "中油上海", "中油港汇",
    "中石油上港", "中油浦东", "中油华鑫", "中油中燃",
]
CHART_COLORS = [
    "4472C4", "ED7D31", "FFC000", "70AD47", "00B0F0",
    "A5A5A5", "5B9BD5", "264478", "C55A11",
]
BLUE = RGBColor(47, 117, 181)
GRID = RGBColor(217, 225, 233)
FONT_SANS = "Noto Sans CJK SC"
FONT_SERIF = "Noto Serif CJK SC"
KEY_ISSUE_CATEGORIES = ["重点商品", "月度盘点", "商品过期", "团购问题"]
CATEGORY_DISPLAY_NAMES = {
    "员工形象及开口服务情况": "员工服务",
    "便利店卫生情况": "便利店卫生",
    "店销商品摆放情况": "商品摆放",
    "仓库管理情况": "仓库管理",
    "销售行为": "销售行为",
    "商品订单、入库、盘点等情况": "商品盘点",
    "台账、报表情况": "台账报表",
    "特殊扣分项": "特殊扣分项",
}


def _unit_name(value):
    text = str(value or "").strip()
    return text[:-2] if text.endswith("片区") else text


def _unit_sort_key(item):
    name = _unit_name(item.get("unit_name"))
    try:
        return UNIT_ORDER.index(name), name
    except ValueError:
        return len(UNIT_ORDER), name


def _format_number(value, digits=0):
    try:
        number = float(value or 0)
    except (TypeError, ValueError):
        number = 0
    return f"{number:.{digits}f}" if digits else str(int(round(number)))


def _shape_text(shape):
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(item.text for item in shape.text_frame.paragraphs).strip()


def _find_text_shape(slide, keyword):
    return next((shape for shape in slide.shapes if keyword in _shape_text(shape)), None)


def _iter_shapes_recursive(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes_recursive(shape.shapes)


def _capture_run_style(text_frame):
    paragraph = next(
        (item for item in text_frame.paragraphs if item.runs),
        text_frame.paragraphs[0],
    )
    run = paragraph.runs[0] if paragraph.runs else None
    font = run.font if run else None
    color = None
    if font is not None:
        try:
            color = font.color.rgb
        except (AttributeError, TypeError, ValueError):
            color = None
    return {
        "alignment": paragraph.alignment,
        "level": paragraph.level,
        "font_name": font.name if font else None,
        "font_size": font.size if font else None,
        "bold": font.bold if font else None,
        "italic": font.italic if font else None,
        "color": color,
    }


def _portable_font_name(source_name):
    value = str(source_name or "").lower()
    if any(keyword in value for keyword in ("仿宋", "小标宋", "宋体", "fangsong", "simsun", "noto serif")):
        return FONT_SERIF
    return FONT_SANS


def _set_run_typeface(run, font_name=None):
    typeface = font_name or _portable_font_name(run.font.name)
    run.font.name = typeface
    run_properties = run._r.get_or_add_rPr()
    for tag_name in ("a:latin", "a:ea", "a:cs"):
        typeface_element = run_properties.find(qn(tag_name))
        if typeface_element is None:
            typeface_element = OxmlElement(tag_name)
            run_properties.append(typeface_element)
        typeface_element.set("typeface", typeface)


def _normalize_text_frame_fonts(text_frame):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            _set_run_typeface(run, _portable_font_name(run.font.name))


def _normalize_shape_fonts(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _normalize_shape_fonts(child)
    if getattr(shape, "has_text_frame", False):
        _normalize_text_frame_fonts(shape.text_frame)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                _normalize_text_frame_fonts(cell.text_frame)
    if getattr(shape, "has_chart", False):
        chart = shape.chart
        if chart.has_title:
            _normalize_text_frame_fonts(chart.chart_title.text_frame)
        try:
            chart.category_axis.tick_labels.font.name = FONT_SANS
            chart.value_axis.tick_labels.font.name = FONT_SANS
        except (AttributeError, ValueError):
            pass
        if chart.has_legend:
            chart.legend.font.name = FONT_SANS
        for plot in chart.plots:
            if getattr(plot, "has_data_labels", False):
                plot.data_labels.font.name = FONT_SANS


def _normalize_presentation_fonts(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            _normalize_shape_fonts(shape)


def _remove_presentation_comments(prs):
    # Drop both legacy and modern comment relationships. Unreferenced parts
    # (including authors) are then omitted by the PPTX package writer.
    for part in list(prs.part.package.iter_parts()):
        for rel in list(part.rels.values()):
            kind = rel.reltype.rsplit("/", 1)[-1].lower()
            if "comment" not in kind and kind not in {"person", "persons"}:
                continue
            element = getattr(part, "_element", None)
            if element is not None:
                for node in list(element.iter()):
                    if node.get(qn("r:id")) == rel.rId:
                        node.getparent().remove(node)
            part.drop_rel(rel.rId)


def _set_summary_bullets(shape, rows, max_font_size=18):
    """Use hanging bullets and separate label/value runs, not template indents."""
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.margin_left = text_frame.margin_right = 0
    text_frame.margin_top = text_frame.margin_bottom = 0
    width = shape.width / 914400
    height = shape.height / 914400
    font_size = max_font_size
    for candidate in range(max_font_size, 9, -1):
        available_units = max(1, (width * 72 - candidate * 1.6) / candidate)
        lines = sum(
            max(1, math.ceil(_visual_text_units(label + value) / available_units))
            for label, value, _color in rows
        )
        if (lines * candidate * 1.35 + (len(rows) - 1) * candidate * 0.38) / 72 <= height:
            font_size = candidate
            break
    else:
        font_size = 10
    for index, (label, value, color) in enumerate(rows):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.clear()
        paragraph.level = 0
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = 1.25
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(font_size * 0.38)
        ppr = paragraph._p.get_or_add_pPr()
        # Remove the source template's tab stops, font bullets and spacing.
        for child in list(ppr):
            if child.tag in {qn("a:tabLst"), qn("a:defRPr")} or "bu" in child.tag.rsplit("}", 1)[-1]:
                ppr.remove(child)
        ppr.set("marL", str(Pt(font_size * 1.55)))
        ppr.set("indent", str(-Pt(font_size * 1.55)))
        bullet_font = OxmlElement("a:buFont")
        bullet_font.set("typeface", FONT_SANS)
        ppr.append(bullet_font)
        bullet = OxmlElement("a:buChar")
        bullet.set("char", "➢")
        ppr.append(bullet)
        for text, bold in ((label, True), (value, False)):
            if not text:
                continue
            run = paragraph.add_run()
            run.text = text
            _set_run_typeface(run, FONT_SERIF)
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.italic = False
            run.font.color.rgb = RGBColor.from_string(color or "111111")


def _set_text_frame(text_frame, value, style=None, font_size=None, bold=None):
    style = style or _capture_run_style(text_frame)
    lines = str(value or "").split("\n")
    text_frame.clear()
    text_frame.word_wrap = True
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = style.get("alignment")
        paragraph.level = style.get("level") or 0
        if not paragraph.runs:
            continue
        run = paragraph.runs[0]
        if style.get("font_name"):
            _set_run_typeface(run, _portable_font_name(style["font_name"]))
        else:
            _set_run_typeface(run, FONT_SANS)
        effective_font_size = font_size or style.get("font_size")
        if isinstance(effective_font_size, (int, float)) and effective_font_size < 1000:
            effective_font_size = Pt(effective_font_size)
        if effective_font_size:
            run.font.size = effective_font_size
        run.font.bold = bold if bold is not None else style.get("bold")
        run.font.italic = style.get("italic")
        if style.get("color"):
            run.font.color.rgb = style["color"]


def _replace_text_runs(shape, replacements):
    if not getattr(shape, "has_text_frame", False):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            value = run.text
            for old, new in replacements.items():
                value = value.replace(old, new)
            run.text = value


def _remove_shape(shape):
    element = shape._element
    element.getparent().remove(element)


def _delete_slide(prs, slide):
    for slide_id in list(prs.slides._sldIdLst):
        if prs.part.related_part(slide_id.rId) is slide.part:
            prs.part.drop_rel(slide_id.rId)
            prs.slides._sldIdLst.remove(slide_id)
            return


def _renumber_slides(prs):
    """Renumber remaining slides sequentially to fill gaps left by deletions.

    python-pptx assigns new slide partnames via ``len(sldIdLst) + 1`` without
    checking existing partnames in the package.  After deleting slides, this
    causes the next added slide to reuse a partname that still belongs to an
    orphaned (but not-yet-garbage-collected) slide part, producing duplicate
    entries in the ZIP and a corrupted PPTX that LibreOffice cannot convert.
    """
    for new_num, slide_id in enumerate(prs.slides._sldIdLst, start=1):
        slide_part = prs.part.related_part(slide_id.rId)
        target = PackURI("/ppt/slides/slide%d.xml" % new_num)
        if slide_part.partname != target:
            slide_part.partname = target


def _move_slide(prs, slide, target_index):
    for slide_id in list(prs.slides._sldIdLst):
        if prs.part.related_part(slide_id.rId) is slide.part:
            prs.slides._sldIdLst.remove(slide_id)
            prs.slides._sldIdLst.insert(target_index, slide_id)
            return


def _set_cell_text(cell, value, font_size=None, bold=None):
    _set_text_frame(
        cell.text_frame,
        value,
        style=_capture_run_style(cell.text_frame),
        font_size=font_size,
        bold=bold,
    )
    cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER


def _resize_table_rows(table, target_row_count):
    table_element = table._tbl
    while len(table.rows) < target_row_count:
        table_element.append(deepcopy(table_element.tr_lst[-1]))
    while len(table.rows) > target_row_count:
        table_element.remove(table_element.tr_lst[-1])


def _fill_table(shape, headers, rows, font_size=11):
    target_height = shape.height
    table = shape.table
    _resize_table_rows(table, len(rows) + 1)
    for column_index, header in enumerate(headers):
        _set_cell_text(table.cell(0, column_index), header, font_size=font_size, bold=True)
    for row_index, values in enumerate(rows, 1):
        for column_index, value in enumerate(values):
            _set_cell_text(
                table.cell(row_index, column_index),
                value,
                font_size=font_size,
                bold=False,
            )
    if len(rows) > 14:
        for table_row in table.rows:
            for cell in table_row.cells:
                cell.margin_top = 0
                cell.margin_bottom = 0
    header_height = min(Inches(0.48), int(target_height * 0.16))
    table.rows[0].height = header_height
    row_height = max(Inches(0.12), int((target_height - header_height) / max(1, len(rows))))
    for row in list(table.rows)[1:]:
        row.height = row_height
    shape.height = target_height


def _visual_text_units(value):
    units = 0.0
    for character in str(value or ""):
        if character == "\n":
            continue
        units += 0.55 if ord(character) < 128 else 1.0
    return units


def _wrap_scope_cell(value, width_inches, font_size):
    # Reserve a little width for font differences between Office and the renderer.
    limit = max(1, (width_inches * 72 - 6) * 0.90 / font_size)
    lines = []
    for paragraph in str(value or "").split("\n"):
        line, units = "", 0.0
        for character in paragraph:
            character_units = _visual_text_units(character)
            if line and (ord(line[-1]) < 128) != (ord(character) < 128):
                # Office inserts spacing at Latin/CJK boundaries, e.g. 站名（12）.
                character_units += 0.3
            if line and units + character_units > limit:
                split_index = line.rfind("、") + 1
                if split_index:
                    lines.append(line[:split_index])
                    line = line[split_index:]
                    units = _visual_text_units(line) + 0.3 * sum(
                        (ord(left) < 128) != (ord(right) < 128)
                        for left, right in zip(line, line[1:])
                    )
                else:
                    lines.append(line)
                    line, units = "", 0.0
            line += character
            units += character_units
        lines.append(line)
    return lines


def _set_scope_text(text_frame, period_text, scope_text, font_size):
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_top = Inches(0.04)
    text_frame.margin_bottom = Inches(0.04)
    text_frame.margin_left = text_frame.margin_right = 0
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    for index, (label, value) in enumerate(
        (("巡检期间：", period_text), ("巡检范围：", scope_text))
    ):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph._p.get_or_add_pPr().clear()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = Pt(font_size * 1.3)
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(font_size * 0.2 if index == 0 else 0)
        label_run = paragraph.add_run()
        label_run.text = label
        label_run.font.size = Pt(font_size)
        label_run.font.bold = True
        _set_run_typeface(label_run, FONT_SERIF)
        value_run = paragraph.add_run()
        value_run.text = str(value or "-")
        value_run.font.size = Pt(font_size)
        value_run.font.bold = False
        _set_run_typeface(value_run, FONT_SERIF)


def _set_chart_fonts(chart, category_size=9):
    try:
        chart.category_axis.tick_labels.font.size = Pt(category_size)
        chart.category_axis.tick_labels.font.name = FONT_SANS
    except (AttributeError, ValueError):
        pass
    try:
        chart.value_axis.tick_labels.font.size = Pt(9)
        chart.value_axis.tick_labels.font.name = FONT_SANS
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = GRID
    except (AttributeError, ValueError):
        pass
    if chart.has_legend:
        chart.legend.font.size = Pt(9)
        chart.legend.font.name = FONT_SANS


def _add_empty_chart_message(slide, box, title):
    text_box = slide.shapes.add_textbox(*box)
    _set_text_frame(
        text_box.text_frame,
        f"{title}\n当前范围暂无数据",
        font_size=Pt(18),
    )
    for paragraph in text_box.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER


def _add_chart_data_table(chart, font_size):
    """Native chart data table: values stay editable with the chart workbook."""
    table = OxmlElement("c:dTable")
    for name in ("showHorzBorder", "showVertBorder", "showOutline", "showKeys"):
        node = OxmlElement(f"c:{name}")
        node.set("val", "1")
        table.append(node)
    properties = OxmlElement("c:spPr")
    line = OxmlElement("a:ln")
    line.set("w", str(Pt(0.5)))
    fill = OxmlElement("a:solidFill")
    color = OxmlElement("a:srgbClr")
    color.set("val", "D9D9D9")
    fill.append(color)
    line.append(fill)
    properties.append(line)
    table.append(properties)
    text = OxmlElement("c:txPr")
    text.append(OxmlElement("a:bodyPr"))
    text.append(OxmlElement("a:lstStyle"))
    paragraph = OxmlElement("a:p")
    ppr = OxmlElement("a:pPr")
    defaults = OxmlElement("a:defRPr")
    defaults.set("sz", str(int(font_size * 100)))
    for tag in ("a:latin", "a:ea", "a:cs"):
        font = OxmlElement(tag)
        font.set("typeface", FONT_SANS)
        defaults.append(font)
    ppr.append(defaults)
    paragraph.append(ppr)
    text.append(paragraph)
    table.append(text)
    chart._chartSpace.chart.plotArea.insert_element_before(table, "c:spPr", "c:extLst")
    chart.has_legend = False
    chart.category_axis.tick_label_position = XL_TICK_LABEL_POSITION.NONE


def _add_column_chart(slide, box, categories, series, title, *, data_table=False):
    categories = [str(item or "-") for item in categories]
    if not categories:
        _add_empty_chart_message(slide, box, title)
        return None
    chart_data = ChartData()
    chart_data.categories = categories
    for name, values, _color in series:
        chart_data.add_series(name, [float(value or 0) for value in values])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        *box,
        chart_data,
    ).chart
    chart.chart_style = 10
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    title_run = chart.chart_title.text_frame.paragraphs[0].runs[0]
    title_run.font.size = Pt(14)
    title_run.font.bold = True
    _set_run_typeface(title_run, FONT_SANS)
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    plot = chart.plots[0]
    plot.gap_width = 65 if len(categories) <= 8 else 35
    plot.has_data_labels = not data_table
    if not data_table:
        plot.data_labels.show_value = True
        plot.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
        plot.data_labels.font.size = Pt(8)
    for index, chart_series in enumerate(chart.series):
        color = series[index][2] if index < len(series) else CHART_COLORS[index % len(CHART_COLORS)]
        chart_series.format.fill.solid()
        chart_series.format.fill.fore_color.rgb = RGBColor.from_string(color)
        chart_series.format.line.color.rgb = RGBColor.from_string(color)
    _set_chart_fonts(chart, category_size=8 if len(categories) > 10 else 9)
    if data_table:
        _add_chart_data_table(chart, 8 if len(categories) > 10 or len(series) > 5 else 9)
    return chart


def _add_pie_chart(slide, box, rows, title):
    rows = [item for item in rows or [] if float(item.get("count") or 0) > 0]
    if not rows:
        _add_empty_chart_message(slide, box, title)
        return None
    chart_data = ChartData()
    chart_data.categories = [str(item.get("name") or "-") for item in rows]
    chart_data.add_series("问题数量", [float(item.get("count") or 0) for item in rows])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, *box, chart_data).chart
    chart.chart_style = 10
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    title_run = chart.chart_title.text_frame.paragraphs[0].runs[0]
    title_run.font.size = Pt(14)
    title_run.font.bold = True
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)
    chart.legend.font.name = FONT_SANS
    plot = chart.plots[0]
    plot.has_data_labels = True
    # New chart labels default to showing values as well as percentages.
    plot.data_labels.show_value = False
    plot.data_labels.show_category_name = False
    plot.data_labels.show_series_name = False
    plot.data_labels.show_legend_key = False
    plot.data_labels.show_percentage = True
    plot.data_labels.position = XL_DATA_LABEL_POSITION.BEST_FIT
    plot.data_labels.font.size = Pt(8)
    for index, point in enumerate(chart.series[0].points):
        color = CHART_COLORS[index % len(CHART_COLORS)]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor.from_string(color)
    return chart


def _set_cell_fill(cell, color):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def _add_styled_table(slide, box, headers, rows, font_size=10, first_column_width=None):
    row_count = max(2, len(rows) + 1)
    table_shape = slide.shapes.add_table(
        row_count,
        len(headers),
        *box,
    )
    table = table_shape.table
    if first_column_width and len(headers) > 1:
        table.columns[0].width = Inches(first_column_width)
        remaining = table_shape.width - table.columns[0].width
        column_width = int(remaining / (len(headers) - 1))
        for column in list(table.columns)[1:]:
            column.width = column_width
    for column_index, header in enumerate(headers):
        cell = table.cell(0, column_index)
        _set_cell_text(cell, header, font_size=font_size, bold=True)
        _set_cell_fill(cell, RGBColor(124, 211, 235))
    normalized_rows = rows or [["暂无数据"] + [""] * (len(headers) - 1)]
    for row_index, values in enumerate(normalized_rows, 1):
        for column_index in range(len(headers)):
            value = values[column_index] if column_index < len(values) else ""
            cell = table.cell(row_index, column_index)
            _set_cell_text(cell, value, font_size=font_size, bold=False)
            _set_cell_fill(
                cell,
                RGBColor(239, 246, 255) if row_index % 2 else RGBColor(255, 255, 255),
            )
    header_height = Inches(0.4)
    table.rows[0].height = header_height
    body_height = max(Inches(0.24), int((table_shape.height - header_height) / max(1, len(normalized_rows))))
    for row in list(table.rows)[1:]:
        row.height = body_height
    return table_shape


def _resolve_image_source(raw_path, storage_root):
    value = str(raw_path or "").strip()
    if not value or not storage_root:
        return None
    clean = value.split("?", 1)[0].replace("\\", "/")
    if "/storage/" in clean:
        clean = clean.split("/storage/", 1)[1]
    clean = clean.lstrip("/")
    root = os.path.abspath(storage_root)
    candidate = os.path.abspath(os.path.join(root, clean))
    try:
        if os.path.commonpath([root, candidate]) != root:
            return None
    except ValueError:
        return None
    return candidate if os.path.isfile(candidate) else None


def _add_picture_contain(slide, raw_path, box, storage_root):
    left, top, width, height = box
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(241, 245, 249)
    frame.line.color.rgb = RGBColor(203, 213, 225)
    source = _resolve_image_source(raw_path, storage_root)
    if not source:
        _set_text_frame(frame.text_frame, "暂无问题照片", font_size=Pt(11))
        for paragraph in frame.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
        frame.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        return None
    try:
        with Image.open(source) as image:
            normalized = ImageOps.exif_transpose(image)
            image_width, image_height = normalized.size
        padding = Inches(0.07)
        inner_width = max(Inches(0.1), width - padding * 2)
        inner_height = max(Inches(0.1), height - padding * 2)
        scale = min(inner_width / max(1, image_width), inner_height / max(1, image_height))
        target_width = int(image_width * scale)
        target_height = int(image_height * scale)
        picture_left = left + int((width - target_width) / 2)
        picture_top = top + int((height - target_height) / 2)
        return slide.shapes.add_picture(
            source,
            picture_left,
            picture_top,
            width=target_width,
            height=target_height,
        )
    except Exception:
        _set_text_frame(frame.text_frame, "照片读取失败", font_size=Pt(11))
        for paragraph in frame.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
        frame.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        return None


def _add_issue_photo_card(slide, issue, box, storage_root, caption_height=0.42):
    left, top, width, height = box
    caption_height_value = Inches(caption_height)
    photo_height = max(Inches(0.3), height - caption_height_value)
    _add_picture_contain(
        slide,
        (issue or {}).get("issue_photo"),
        (left, top, width, photo_height),
        storage_root,
    )
    caption = slide.shapes.add_textbox(
        left,
        top + photo_height,
        width,
        caption_height_value,
    )
    has_issue = bool(issue and (issue.get("issue_id") or issue.get("description")))
    station_name = str((issue or {}).get("station_name") or "")
    description = str((issue or {}).get("description") or "").strip()
    if len(description) > 34:
        description = f"{description[:34]}…"
    _set_text_frame(
        caption.text_frame,
        f"{station_name}：{description}" if has_issue else "暂无更多典型问题",
        font_size=Pt(9),
        bold=False,
    )
    for paragraph in caption.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
    return caption


def _feature_distribution(issues, feature_groups):
    rows = []
    for name, keywords in feature_groups:
        count = sum(
            1 for issue in issues or []
            if any(keyword in str(issue.get("description") or "") for keyword in keywords)
        )
        if count:
            rows.append({"name": name, "count": count})
    return rows


def _find_key_detail(report, category_name):
    return next(
        (
            item for item in (report.get("key_issue_summary") or {}).get("details") or []
            if item.get("name") == category_name
        ),
        {"name": category_name, "count": 0, "percentage_of_all": 0, "relationship": [], "issues": []},
    )


def _remove_large_visuals(slide):
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.CHART:
            _remove_shape(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.height > Inches(0.3):
            _remove_shape(shape)


def _edit_cover(slide, report):
    month_number = int(str(report.get("month") or "2000-01").split("-")[1])
    title_shape = _find_text_shape(slide, "盛大业务督导中心")
    if title_shape:
        _set_text_frame(
            title_shape.text_frame,
            f"盛大业务督导中心\n非油业务{month_number}月巡检报告",
            style=_capture_run_style(title_shape.text_frame),
        )
    name_shape = _find_text_shape(slide, "王昕怡")
    if name_shape:
        _set_text_frame(
            name_shape.text_frame,
            "XXX    XXX    XXX",
            style=_capture_run_style(name_shape.text_frame),
        )


def _edit_rectification_slide(slide, report):
    previous = report.get("previous_month_rectification") or {}
    narrative_shape = _find_text_shape(slide, "各片区整改情况")
    if narrative_shape:
        _set_text_frame(
            narrative_shape.text_frame,
            previous.get("narrative") or "当前范围暂无上期整改数据。",
        )
    _remove_large_visuals(slide)
    units = sorted(previous.get("units") or [], key=_unit_sort_key)
    _add_column_chart(
        slide,
        (Inches(1.25), Inches(2.08), Inches(10.85), Inches(4.55)),
        [_unit_name(item.get("unit_name")) for item in units],
        [
            ("全部问题", [item.get("total_count") for item in units], "4472C4"),
            ("待验收", [item.get("pending_acceptance_count") for item in units], "C0504D"),
            ("待整改", [item.get("pending_rectification_count") for item in units], "70AD47"),
        ],
        "各片区待验收与待整改问题数量分布",
        data_table=True,
    )


def _scope_detail_rows(report):
    units = sorted(report.get("units") or [], key=_unit_sort_key)
    rows = []
    for index, item in enumerate(units, 1):
        station_text = "、".join(
            f"{row.get('station_name')}（{row.get('issue_count') or 0}）"
            for row in item.get("station_issue_rows") or []
        )
        rows.append(
            [
                str(index),
                _unit_name(item.get("unit_name")),
                str(item.get("station_count") or 0),
                str(item.get("issue_count") or 0),
                _format_number(item.get("average_issue_count"), 1),
                station_text,
            ]
        )
    return rows or [["-", "暂无数据", "0", "0", "0.0", "-"]]


def _scope_detail_layout(lines, width, column_widths, rows, available_height):
    candidates = []
    for table_font in range(12, 5, -1):
        wrapped_rows = [
            [_wrap_scope_cell(value, column_widths[index], table_font) for index, value in enumerate(row)]
            for row in rows
        ]
        row_heights = [
            (max(len(cell) for cell in row) * table_font * 1.25 + 2) / 72
            for row in wrapped_rows
        ]
        for text_font in range(20, 11, -1):
            line_count = sum(len(_wrap_scope_cell(line, width, text_font)) for line in lines)
            text_height = (line_count * text_font * 1.3 + text_font * 0.2) / 72 + 0.08
            if text_height + 0.12 + sum(row_heights) > available_height:
                continue
            candidates.append({
                "text_font": text_font,
                "text_height": text_height,
                "table_font": table_font,
                "wrapped_rows": wrapped_rows,
                "row_heights": row_heights,
                "score": table_font * 5 + text_font * 0.3,
            })
    if not candidates:
        # Never produce a clipped slide or silently discard station details.
        raise ValueError("巡检范围明细超出单页可读容量，请缩小报告日期范围后重新生成。")
    return max(candidates, key=lambda candidate: candidate["score"])


def _edit_scope_slide(slide, report, detail_slide):
    text_shape = _find_text_shape(slide, "巡检期间")
    source_table = next(shape for shape in detail_slide.shapes if shape.has_table)
    if text_shape is None:
        raise ValueError("非油报告巡检范围模板缺少说明文本框。")
    for old_table in list(slide.shapes):
        if old_table.has_table:
            _remove_shape(old_table)
    table_element = deepcopy(source_table._element)
    identity = table_element.find(".//" + qn("p:cNvPr"))
    identity.set("id", str(slide.shapes._next_shape_id))
    identity.set("name", "巡检范围及各站问题明细表")
    slide.shapes._spTree.insert_element_before(table_element, "p:extLst")
    table_shape = slide.shapes[-1]
    table = table_shape.table

    text_shape.top = Inches(1.20)
    table_shape.left = text_shape.left
    table_shape.width = text_shape.width
    original_width = sum(column.width for column in table.columns)
    for column in table.columns:
        column.width = int(column.width / original_width * text_shape.width)
    headers = ["序号", "所属片区", "站点数量", "片区问题总项", "站平均问题数", "站点问题数"]
    rows = [headers, *_scope_detail_rows(report)]
    period_text = report.get("period_text") or "-"
    scope_text = report.get("scope_text") or "-"
    content_bottom = 7.27
    available_height = content_bottom - text_shape.top / 914400
    layout = _scope_detail_layout(
        ("巡检期间：" + period_text, "巡检范围：" + scope_text),
        text_shape.width / 914400,
        [column.width / 914400 for column in table.columns],
        rows,
        available_height,
    )
    text_shape.height = Inches(layout["text_height"])
    _set_scope_text(text_shape.text_frame, period_text, scope_text, layout["text_font"])
    table_shape.top = text_shape.top + text_shape.height + Inches(0.12)
    _resize_table_rows(table, len(rows))
    # Give wrapped station lists more height rather than squeezing every row equally.
    available_table_height = content_bottom - table_shape.top / 914400
    expansion = min(1.5, available_table_height / sum(layout["row_heights"]))
    for row_index, wrapped_row in enumerate(layout["wrapped_rows"]):
        table.rows[row_index].height = Inches(layout["row_heights"][row_index] * expansion)
        for column_index, wrapped in enumerate(wrapped_row):
            cell = table.cell(row_index, column_index)
            cell.margin_left = cell.margin_right = Pt(3)
            cell.margin_top = cell.margin_bottom = Pt(1)
            _set_cell_text(cell, "\n".join(wrapped), font_size=layout["table_font"], bold=row_index == 0)
            for paragraph in cell.text_frame.paragraphs:
                paragraph._p.get_or_add_pPr().clear()
                for end_style in paragraph._p.findall(qn("a:endParaRPr")):
                    paragraph._p.remove(end_style)
                paragraph.alignment = PP_ALIGN.LEFT if column_index == 5 and row_index else PP_ALIGN.CENTER
                paragraph.line_spacing = Pt(layout["table_font"] * 1.25)
                paragraph.space_before = paragraph.space_after = Pt(0)
                for run in paragraph.runs:
                    _set_run_typeface(run, FONT_SANS)
    table_shape.height = sum(row.height for row in table.rows)


def _edit_overview_slide(slide, report):
    _remove_large_visuals(slide)
    title_id = slide.shapes.title.shape_id if slide.shapes.title else None
    summary_shape = next(
        (
            shape for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and shape.shape_id != title_id
            and shape.top > Inches(4.8)
        ),
        None,
    )
    units = sorted(report.get("units") or [], key=_unit_sort_key)
    summary = report.get("summary") or {}
    categories = sorted(
        [item for item in report.get("category_distribution") or [] if item.get("count")],
        key=lambda item: -item["count"],
    )
    category_total = sum(item["count"] for item in categories)
    concentration = "当前范围暂无非油现场检查问题。"
    if categories:
        concentration = (
            f"非油现场检查问题主要集中在{categories[0]['name']}，"
            f"占比{categories[0]['count'] / category_total * 100:.1f}%"
        )
        if len(categories) > 1:
            concentration += (
                f"，其次为{categories[1]['name']}（"
                f"{categories[1]['count'] / category_total * 100:.1f}%）"
            )
        concentration += "。"
    station_count = summary.get("station_count", sum(item.get("station_count") or 0 for item in units))
    total = summary.get("total_issue_count") or 0
    month = int(str(report.get("month") or "2000-01").split("-")[1])
    if summary_shape:
        summary_shape.left, summary_shape.top = Inches(0.55), Inches(5.40)
        summary_shape.width, summary_shape.height = Inches(12.25), Inches(1.88)
        _set_summary_bullets(summary_shape, [
            ("片区（分公司）总数：", f"{sum(item.get('unit_type') == 'region' for item in units)}个管理片区和{sum(item.get('unit_type') == 'holding' for item in units)}个合资公司", None),
            ("问题总计：", f"{total}项", None),
            ("站平均问题数：", f"{total / station_count if station_count else 0:.1f}项/站", None),
            ("", concentration, "C00000"),
            ("具体问题详见附件：", f"《{month}月非油检查问题清单》", None),
        ], max_font_size=16)
    _add_column_chart(
        slide,
        (Inches(0.25), Inches(1.05), Inches(7.55), Inches(4.05)),
        [_unit_name(item.get("unit_name")) for item in units],
        [
            ("单站平均问题数", [item.get("average_issue_count") for item in units], "0070C0"),
            ("检查站点数", [item.get("station_count") for item in units], "C0504D"),
            ("问题总数", [item.get("issue_count") for item in units], "92D050"),
        ],
        "各片区非油问题数量汇总",
        data_table=True,
    )
    _add_pie_chart(
        slide,
        (Inches(7.85), Inches(1.08), Inches(5.15), Inches(4.05)),
        report.get("category_distribution") or [],
        "非油检查问题分布",
    )


def _copy_text_shape(source_shape, destination_slide):
    element = deepcopy(source_shape._element)
    destination_slide.shapes._spTree.insert_element_before(element, "p:extLst")
    return destination_slide.shapes[-1]


def _add_unit_slide(prs, source_slide, unit):
    slide = prs.slides.add_slide(source_slide.slide_layout)
    title_shape = slide.shapes.title
    title_style = _capture_run_style(source_slide.shapes.title.text_frame)
    unit_label = _unit_name(unit.get("unit_name"))
    if unit.get("unit_type") == "region":
        unit_label += "片区"
    if title_shape:
        _set_text_frame(
            title_shape.text_frame,
            f"二、片区问题汇总——{unit_label}",
            style=title_style,
        )
    source_summary = _find_text_shape(source_slide, "涉及站点数")
    if source_summary is None:
        raise ValueError("非油报告片区模板缺少汇总文本框。")
    summary_shape = _copy_text_shape(source_summary, slide)
    station_rows = unit.get("station_issue_rows") or []
    station_names = "、".join(unit.get("station_names") or [item.get("station_name") or "" for item in station_rows])
    summary_shape.left = Inches(0.75)
    summary_shape.top = Inches(5.35)
    summary_shape.width = Inches(11.9)
    summary_shape.height = Inches(1.75)
    _set_summary_bullets(summary_shape, [
        ("涉及站点数：", f"{unit.get('station_count') or 0}座（{station_names}）", None),
        ("问题总计：", f"{unit.get('issue_count') or 0}项（占本次全盘通报问题总量的{float(unit.get('percentage') or 0):.1f}%）", None),
        ("站平均问题数：", f"{float(unit.get('average_issue_count') or 0):.1f}项/站", None),
    ], max_font_size=20)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        Inches(5.22),
        prs.slide_width,
        Inches(0.06),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    categories = [item.get("station_name") or "-" for item in station_rows]
    category_names = [
        item.get("name")
        for item in unit.get("category_distribution") or []
        if item.get("count")
    ]
    series = [
        (
            CATEGORY_DISPLAY_NAMES.get(category_name, category_name),
            [
                (station.get("category_counts") or {}).get(category_name, 0)
                for station in station_rows
            ],
            CHART_COLORS[index % len(CHART_COLORS)],
        )
        for index, category_name in enumerate(category_names)
    ]
    if not series:
        series = [
            (
                "问题数量",
                [station.get("issue_count") or 0 for station in station_rows],
                "4472C4",
            )
        ]
    _add_column_chart(
        slide,
        (Inches(0.45), Inches(1.15), Inches(7.0), Inches(3.85)),
        categories,
        series,
        f"{_unit_name(unit.get('unit_name'))}各站问题数量",
        data_table=True,
    )
    _add_pie_chart(
        slide,
        (Inches(7.55), Inches(1.15), Inches(5.25), Inches(3.85)),
        unit.get("category_distribution") or [],
        f"{_unit_name(unit.get('unit_name'))}问题分布",
    )
    return slide


def _edit_analysis_overview(slide, report):
    summary = report.get("summary") or {}
    total = int(summary.get("total_issue_count") or 0)
    selected = (report.get("key_issue_summary") or {}).get("selected_count")
    if selected is None:
        selected = report.get("key_issue_count") or 0
    selected = int(selected)
    percentage = f"{selected / total * 100 if total else 0:.1f}".rstrip("0").rstrip(".")

    # Cards are nested groups. Replace only their numeric runs so their colors,
    # labels and line breaks survive; retain the fixed methodology verbatim.
    for label, value in (
        ("问题总数", str(total)),
        ("重点问题数", str(selected)),
        ("重点问题占比", f"{percentage}%"),
    ):
        shape = next(
            (item for item in _iter_shapes_recursive(slide.shapes) if label in _shape_text(item)),
            None,
        )
        if shape is None:
            raise ValueError(f"非油报告检查总体情况模板缺少统计项：{label}")
        value_run = next(
            (
                run for paragraph in shape.text_frame.paragraphs for run in paragraph.runs
                if re.fullmatch(r"\s*\d+(?:\.\d+)?%?\s*", run.text)
            ),
            None,
        )
        if value_run is None:
            raise ValueError(f"非油报告检查总体情况模板缺少数值：{label}")
        value_run.text = re.sub(r"\d+(?:\.\d+)?%?", value, value_run.text, count=1)


def _edit_key_issue_overview_slide(slide, report):
    for shape in list(slide.shapes):
        if shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART}:
            _remove_shape(shape)
    for category_name in KEY_ISSUE_CATEGORIES:
        label_shape = next(
            (
                shape for shape in _iter_shapes_recursive(slide.shapes)
                if _shape_text(shape) == category_name
            ),
            None,
        )
        if not label_shape:
            continue
        label_shape.text_frame.margin_left = 0
        label_shape.text_frame.margin_right = 0
        label_shape.text_frame.margin_top = 0
        label_shape.text_frame.margin_bottom = 0
        _set_text_frame(
            label_shape.text_frame,
            category_name,
            font_size=Pt(9),
            bold=True,
        )
        label_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in label_shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
    _add_pie_chart(
        slide,
        (Inches(8.82), Inches(2.82), Inches(4.2), Inches(3.55)),
        (report.get("key_issue_summary") or {}).get("distribution") or [],
        "重点问题内部占比",
    )


def _edit_key_issue_relationship_slide(slide, report):
    summary = report.get("key_issue_summary") or {}
    details = [item for item in summary.get("details") or [] if item.get("count")]
    narrative_shape = _find_text_shape(slide, "根据问题定义")
    narrative_lines = ["根据问题定义，将重点问题与检查环节做关联性分析，可以得出："]
    for detail in details:
        relationships = detail.get("relationship") or []
        top_rows = relationships[:2]
        if not top_rows:
            relationship_text = "当前未形成明确的业务环节分布"
        else:
            relationship_text = "、".join(
                f"{item.get('name')}（{_format_number(item.get('percentage'), 1)}%）"
                for item in top_rows
            )
        narrative_lines.append(f"·{detail.get('name')}发生在{relationship_text}。")
    if len(narrative_lines) == 1:
        narrative_lines.append("·当前范围暂无符合四类定义的重点问题。")
    if narrative_shape:
        narrative_shape.height = Inches(2.1)
        _set_text_frame(narrative_shape.text_frame, "\n".join(narrative_lines), font_size=Pt(16))

    for shape in list(slide.shapes):
        if getattr(shape, "has_table", False):
            _remove_shape(shape)
    active_names = []
    for source_name, display_name in CATEGORY_DISPLAY_NAMES.items():
        if any(
            item.get("source_name") == source_name
            for detail in details
            for item in detail.get("relationship") or []
        ):
            active_names.append(display_name)
    if not active_names:
        active_names = ["暂无环节数据"]
    rows = []
    for category_name in KEY_ISSUE_CATEGORIES:
        detail = next((item for item in details if item.get("name") == category_name), None)
        relationship_map = {
            item.get("name"): item for item in (detail or {}).get("relationship") or []
        }
        rows.append(
            [category_name]
            + [
                (
                    f"{_format_number(relationship_map[name].get('percentage'), 1)}%"
                    if name in relationship_map else ""
                )
                for name in active_names
            ]
        )
    _add_styled_table(
        slide,
        (Inches(0.88), Inches(4.45), Inches(11.55), Inches(2.25)),
        ["重点问题"] + active_names,
        rows,
        font_size=10 if len(active_names) > 6 else 12,
        first_column_width=1.7,
    )


def _edit_key_product_overview_slide(slide, report):
    detail = _find_key_detail(report, "重点商品")
    issues = detail.get("issues") or []
    relationship = detail.get("relationship") or []
    brand_distribution = _feature_distribution(
        issues,
        [
            ("中华烟", ("中华", "软中", "硬中")),
            ("五粮液", ("五粮液",)),
            ("茅台", ("茅台", "贵州茅台")),
        ],
    )
    manifestation_distribution = _feature_distribution(
        issues,
        [
            ("盘亏", ("盘亏",)),
            ("盘盈", ("盘盈",)),
            ("未过机", ("未过机", "不过机")),
            ("未出样", ("未出样", "未展示")),
            ("价签", ("价签", "标价")),
        ],
    )
    narrative_shape = _find_text_shape(slide, "61%") or _find_text_shape(slide, "重点问题涉及")
    relationship_text = "、".join(
        f"{item.get('name')}（{_format_number(item.get('percentage'), 1)}%）"
        for item in relationship[:3]
    ) or "暂无明确环节分布"
    brand_text = "、".join(
        f"{item.get('name')}（{item.get('count')}项）" for item in brand_distribution
    ) or "暂无可统计品牌"
    manifestation_text = "、".join(
        f"{item.get('name')}（{item.get('count')}项）" for item in manifestation_distribution
    ) or "暂无可统计表现"
    if narrative_shape:
        _set_text_frame(
            narrative_shape.text_frame,
            (
                f"重点商品问题共{detail.get('count') or 0}项，占全部问题"
                f"{_format_number(detail.get('percentage_of_all'), 1)}%\n"
                f"主要发生于{relationship_text}\n"
                f"涉及商品：{brand_text}\n"
                f"具体表现：{manifestation_text}"
            ),
            font_size=Pt(15),
        )
    for shape in list(slide.shapes):
        if shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART}:
            _remove_shape(shape)
    _add_pie_chart(
        slide,
        (Inches(0.55), Inches(3.82), Inches(4.0), Inches(3.05)),
        relationship,
        "发生环节",
    )
    _add_pie_chart(
        slide,
        (Inches(4.65), Inches(3.82), Inches(4.0), Inches(3.05)),
        brand_distribution,
        "涉及商品",
    )
    _add_pie_chart(
        slide,
        (Inches(8.75), Inches(3.82), Inches(4.0), Inches(3.05)),
        manifestation_distribution,
        "典型问题分布",
    )


def _fill_key_issue_table_slide(slide, issues, title_text, summary_text=None):
    bar_shape = next(
        (shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.top > Inches(0.8) and shape.top < Inches(1.6)),
        None,
    )
    if bar_shape:
        _set_text_frame(bar_shape.text_frame, title_text)
    if summary_text:
        summary_shape = next(
            (
                shape for shape in slide.shapes
                if getattr(shape, "has_text_frame", False)
                and shape._element is not bar_shape._element
                and shape.top > Inches(1.25)
                and shape.top < Inches(2.2)
                and not getattr(shape, "has_table", False)
            ),
            None,
        )
        if summary_shape:
            _set_text_frame(summary_shape.text_frame, summary_text, font_size=Pt(15))
    table_shape = next(shape for shape in slide.shapes if getattr(shape, "has_table", False))
    rows = [
        [
            issue.get("source_project") or "-",
            issue.get("station_name") or "-",
            issue.get("description") or "-",
        ]
        for issue in issues
    ]
    _fill_table(
        table_shape,
        ["检查项目", "站点", "问题描述"],
        rows or [["-", "-", "当前范围暂无该类重点问题"]],
        font_size=8 if len(rows) > 8 else 10,
    )


def _edit_key_product_detail_slides(slide_one, slide_two, report):
    detail = _find_key_detail(report, "重点商品")
    issues = detail.get("issues") or []
    _fill_key_issue_table_slide(
        slide_one,
        issues[:10],
        "5.1 重点商品典型问题",
        "重点商品问题主要涉及账实、陈列、销售过机和价签管理。",
    )
    _fill_key_issue_table_slide(
        slide_two,
        issues[10:20],
        "5.2 重点商品典型问题（续）",
    )


def _edit_monthly_inventory_slide(slide, report):
    detail = _find_key_detail(report, "月度盘点")
    _fill_key_issue_table_slide(
        slide,
        (detail.get("issues") or [])[:8],
        "6. 月度盘点问题",
        (
            f"共{detail.get('count') or 0}项，占全部问题"
            f"{_format_number(detail.get('percentage_of_all'), 1)}%\n"
            "重点核实交接班盘点记录、盘点覆盖率及签字完整性。"
        ),
    )


def _edit_group_purchase_slide(slide, report, storage_root):
    detail = _find_key_detail(report, "团购问题")
    issues = detail.get("issues") or []
    summary_shape = _find_text_shape(slide, "共2起") or _find_text_shape(slide, "占全部问题")
    if summary_shape:
        _set_text_frame(
            summary_shape.text_frame,
            f"共{detail.get('count') or 0}项，占全部问题{_format_number(detail.get('percentage_of_all'), 1)}%",
            font_size=Pt(16),
        )
    table_shape = next(shape for shape in slide.shapes if getattr(shape, "has_table", False))
    rows = [
        [
            issue.get("source_project") or "团购合规",
            f"{issue.get('station_name') or '-'}：{issue.get('description') or '-'}",
        ]
        for issue in issues[:4]
    ]
    _fill_table(
        table_shape,
        ["典型问题", "问题详情"],
        rows or [["暂无", "当前范围暂无团购重点问题"]],
        font_size=9,
    )
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (
            getattr(shape, "has_text_frame", False)
            and shape.top > Inches(2.2)
            and shape._element is not table_shape._element
        ):
            _remove_shape(shape)
    for index in range(2):
        issue = issues[index] if index < len(issues) else {}
        _add_issue_photo_card(
            slide,
            issue,
            (Inches(8.0), Inches(2.15 + index * 2.45), Inches(4.45), Inches(2.25)),
            storage_root,
        )


def _edit_expired_product_slide(slide, report, storage_root):
    detail = _find_key_detail(report, "商品过期")
    issues = detail.get("issues") or []
    summary_shape = _find_text_shape(slide, "共7起") or _find_text_shape(slide, "过期商品")
    if summary_shape:
        _set_text_frame(
            summary_shape.text_frame,
            (
                f"共{detail.get('count') or 0}项，占全部问题"
                f"{_format_number(detail.get('percentage_of_all'), 1)}%\n"
                "重点关注过期商品的检查、登记、下架和销毁处理。"
            ),
            font_size=Pt(16),
        )
    preserved_elements = {
        shape._element for shape in list(slide.shapes)[:3]
    }
    for shape in list(slide.shapes):
        if shape._element not in preserved_elements:
            _remove_shape(shape)
    boxes = [
        (Inches(0.75), Inches(2.82), Inches(5.8), Inches(1.82)),
        (Inches(6.78), Inches(2.82), Inches(5.8), Inches(1.82)),
        (Inches(0.75), Inches(4.92), Inches(5.8), Inches(1.82)),
        (Inches(6.78), Inches(4.92), Inches(5.8), Inches(1.82)),
    ]
    for index, box in enumerate(boxes):
        issue = issues[index] if index < len(issues) else {}
        _add_issue_photo_card(slide, issue, box, storage_root)


def _edit_analysis_distribution_slide(slide, report):
    for shape in list(slide.shapes):
        if shape.shape_type in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART}:
            _remove_shape(shape)
    rows = [
        item for item in report.get("category_distribution") or []
        if item.get("count") and item.get("name") in CATEGORY_DISPLAY_NAMES
    ]
    _add_column_chart(
        slide,
        (Inches(0.85), Inches(2.25), Inches(11.7), Inches(4.35)),
        [CATEGORY_DISPLAY_NAMES.get(item.get("name"), item.get("name")) for item in rows],
        [("问题数量", [item.get("count") for item in rows], "4472C4")],
        "非油检查问题分布",
    )


def _edit_analysis_method_slide(slide, report):
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            _remove_shape(shape)
    features = []
    for detail in report.get("category_details") or []:
        for item in detail.get("features") or []:
            features.append((item.get("name") or detail.get("name"), int(item.get("count") or 0)))
    features.sort(key=lambda item: (-item[1], item[0]))
    if not features:
        features = [("暂无高频特征", 1)]
    for index, (name, count) in enumerate(features[:12]):
        column = index % 4
        row = index // 4
        width = 2.65
        left = 1.1 + column * 3.0
        top = 3.18 + row * 0.92
        tag = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(0.65),
        )
        color = RGBColor.from_string(CHART_COLORS[index % len(CHART_COLORS)])
        tag.fill.solid()
        tag.fill.fore_color.rgb = color
        tag.line.fill.background()
        _set_text_frame(tag.text_frame, f"{name}  {count}", font_size=Pt(13), bold=True)
        tag.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in tag.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(255, 255, 255)


def _edit_category_detail_slide(slide, detail, display_index, total_issue_count, storage_root):
    shapes = list(slide.shapes)
    title_shape = shapes[0] if shapes else None
    bar_shape = shapes[1] if len(shapes) > 1 else None
    summary_shape = shapes[2] if len(shapes) > 2 else None
    if title_shape and getattr(title_shape, "has_text_frame", False):
        _set_text_frame(title_shape.text_frame, "四、具体问题分析")
    if bar_shape and getattr(bar_shape, "has_text_frame", False):
        _set_text_frame(bar_shape.text_frame, f"{display_index}. {detail.get('name')}")
    if summary_shape and getattr(summary_shape, "has_text_frame", False):
        _set_text_frame(
            summary_shape.text_frame,
            (
                f"共{detail.get('count') or 0}项，占全部问题"
                f"{_format_number(detail.get('percentage'), 1)}%，典型问题如下："
            ),
            font_size=Pt(15),
        )
    for shape in shapes[3:]:
        _remove_shape(shape)

    feature_rows = [
        [item.get("name") or "-", str(item.get("count") or 0), item.get("example") or "-"]
        for item in detail.get("features") or []
    ]
    _add_styled_table(
        slide,
        (Inches(0.7), Inches(2.82), Inches(5.85), Inches(2.12)),
        ["高频特征", "数量", "典型表现"],
        feature_rows,
        font_size=9,
        first_column_width=1.35,
    )
    text_box = slide.shapes.add_textbox(
        Inches(0.78),
        Inches(5.1),
        Inches(5.72),
        Inches(1.32),
    )
    descriptions = [
        f"·{issue.get('station_name') or '-'}：{issue.get('description') or '-'}"
        for issue in (detail.get("issues") or [])[:3]
    ]
    _set_text_frame(
        text_box.text_frame,
        "\n".join(descriptions) if descriptions else "当前类别暂无可展示的典型问题。",
        font_size=Pt(11),
    )
    photo_boxes = [
        (Inches(6.85), Inches(2.82), Inches(2.8), Inches(1.82)),
        (Inches(9.82), Inches(2.82), Inches(2.8), Inches(1.82)),
        (Inches(6.85), Inches(4.92), Inches(2.8), Inches(1.82)),
        (Inches(9.82), Inches(4.92), Inches(2.8), Inches(1.82)),
    ]
    issues = detail.get("issues") or []
    for index, box in enumerate(photo_boxes):
        issue = issues[index] if index < len(issues) else {}
        _add_issue_photo_card(slide, issue, box, storage_root, caption_height=0.46)


def _add_ai_badge(slide):
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(11.45),
        Inches(1.12),
        Inches(1.12),
        Inches(0.34),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(219, 234, 254)
    badge.line.color.rgb = RGBColor(147, 197, 253)
    _set_text_frame(badge.text_frame, "AI 生成", font_size=Pt(9), bold=True)
    for paragraph in badge.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.color.rgb = RGBColor(29, 78, 216)


def _edit_ai_analysis_slide(slide, report):
    deep_analysis = report.get("deep_analysis") or {}
    core_findings = (deep_analysis.get("core_findings") or [])[:3]
    attribution = (deep_analysis.get("attribution_analysis") or [])[:3]
    core_titles = sorted(
        [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and Inches(2.35) <= shape.top <= Inches(2.8) and shape.height < Inches(0.65)],
        key=lambda shape: shape.left,
    )
    core_contents = sorted(
        [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and Inches(2.75) < shape.top < Inches(4.3) and shape.height > Inches(0.7)],
        key=lambda shape: shape.left,
    )
    attribution_titles = sorted(
        [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and Inches(5.0) <= shape.top <= Inches(5.75) and shape.height < Inches(0.65)],
        key=lambda shape: shape.left,
    )
    attribution_contents = sorted(
        [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.top > Inches(5.7) and shape.height > Inches(0.7)],
        key=lambda shape: shape.left,
    )
    for shapes, items, content_key in (
        (core_titles, core_findings, "title"),
        (core_contents, core_findings, "content"),
        (attribution_titles, attribution, "title"),
        (attribution_contents, attribution, "content"),
    ):
        for index, shape in enumerate(shapes[:3]):
            item = items[index] if index < len(items) else {}
            _set_text_frame(
                shape.text_frame,
                item.get(content_key) or "暂无可用分析",
                font_size=Pt(13 if content_key == "content" else 15),
                bold=content_key == "title",
            )
    _add_ai_badge(slide)


def _edit_ai_actions_slide(slide, report):
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            _remove_shape(shape)
    deep_analysis = report.get("deep_analysis") or {}
    items = (deep_analysis.get("action_priorities") or deep_analysis.get("improvement_suggestions") or [])[:3]
    while len(items) < 3:
        items.append({"title": "待补充", "content": "当前数据不足，请结合后续检查情况完善行动建议。"})
    for index, item in enumerate(items):
        left = 0.72 + index * 4.15
        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(1.75),
            Inches(3.78),
            Inches(4.85),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
        panel.line.color.rgb = RGBColor(226, 232, 240)
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left),
            Inches(1.75),
            Inches(3.78),
            Inches(0.08),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(220, 38, 38)
        accent.line.fill.background()
        number_box = slide.shapes.add_textbox(
            Inches(left + 0.18),
            Inches(2.0),
            Inches(0.55),
            Inches(0.45),
        )
        _set_text_frame(number_box.text_frame, f"0{index + 1}", font_size=Pt(18), bold=True)
        title_box = slide.shapes.add_textbox(
            Inches(left + 0.72),
            Inches(1.98),
            Inches(2.82),
            Inches(0.62),
        )
        _set_text_frame(title_box.text_frame, item.get("title") or "改善建议", font_size=Pt(15), bold=True)
        content_box = slide.shapes.add_textbox(
            Inches(left + 0.25),
            Inches(2.82),
            Inches(3.28),
            Inches(3.2),
        )
        _set_text_frame(content_box.text_frame, item.get("content") or "-", font_size=Pt(13))
    _add_ai_badge(slide)


def _render_presentation_preview(pptx_path, output_dir):
    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    pdftoppm = shutil.which("pdftoppm")
    if not soffice or not pdftoppm:
        raise RuntimeError("服务器缺少 LibreOffice 或 Poppler，无法生成PPT同源预览。")

    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    for stale in output_dir.glob("slide-*.jpg"):
        stale.unlink()

    with tempfile.TemporaryDirectory(prefix="non-oil-ppt-render-") as temp_dir:
        temp_path = Path(temp_dir)
        profile_uri = (temp_path / f"lo-profile-{uuid.uuid4().hex}").as_uri()
        subprocess.run(
            [
                soffice,
                "--headless",
                f"-env:UserInstallation={profile_uri}",
                "--convert-to",
                "pdf",
                "--outdir",
                str(temp_path),
                str(pptx_path),
            ],
            check=True,
            capture_output=True,
            text=True,
            timeout=180,
        )
        pdf_path = temp_path / f"{Path(pptx_path).stem}.pdf"
        if not pdf_path.is_file():
            raise RuntimeError("PPT预览转换失败，未生成PDF中间文件。")
        prefix = temp_path / "page"
        subprocess.run(
            [
                pdftoppm,
                "-jpeg",
                "-scale-to-x",
                str(CANVAS_SIZE[0]),
                "-scale-to-y",
                str(CANVAS_SIZE[1]),
                "-jpegopt",
                "quality=95",
                str(pdf_path),
                str(prefix),
            ],
            check=True,
            capture_output=True,
            text=True,
            timeout=180,
        )
        rendered_pages = sorted(
            temp_path.glob("page-*.jpg"),
            key=lambda path: int(re.search(r"(\d+)$", path.stem).group(1)),
        )
        if not rendered_pages:
            raise RuntimeError("PPT预览转换失败，未生成幻灯片图片。")
        slide_files = []
        for index, source in enumerate(rendered_pages, 1):
            destination = output_dir / f"slide-{index:02d}.jpg"
            shutil.copy2(source, destination)
            slide_files.append(str(destination))
        return slide_files


def build_non_oil_template_presentation(
    report,
    output_dir,
    output_path,
    storage_root=None,
):
    if not TEMPLATE_FILE.is_file():
        raise FileNotFoundError("非油检查报告PPT模板不存在。")
    output_dir = Path(output_dir)
    output_path = Path(output_path)
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation(TEMPLATE_FILE)
    if len(prs.slides) < 45:
        raise ValueError("非油检查报告PPT模板页数异常。")

    _edit_cover(prs.slides[0], report)
    _edit_rectification_slide(prs.slides[3], report)
    detail_slide = prs.slides[5]
    _edit_scope_slide(prs.slides[4], report, detail_slide)
    _edit_overview_slide(prs.slides[8], report)
    _edit_analysis_overview(prs.slides[21], report)
    _edit_key_issue_overview_slide(prs.slides[22], report)
    _edit_key_issue_relationship_slide(prs.slides[23], report)
    _edit_key_product_overview_slide(prs.slides[24], report)
    _edit_key_product_detail_slides(prs.slides[25], prs.slides[26], report)
    _edit_monthly_inventory_slide(prs.slides[27], report)
    _edit_group_purchase_slide(prs.slides[28], report, storage_root)
    _edit_expired_product_slide(prs.slides[29], report, storage_root)
    _edit_analysis_distribution_slide(prs.slides[31], report)
    _edit_analysis_method_slide(prs.slides[32], report)
    _edit_ai_analysis_slide(prs.slides[42], report)
    _edit_ai_actions_slide(prs.slides[43], report)

    category_slides = [prs.slides[index] for index in range(33, 41)]
    category_details = (report.get("category_details") or [])[: len(category_slides)]
    for index, detail in enumerate(category_details):
        _edit_category_detail_slide(
            category_slides[index],
            detail,
            index + 3,
            (report.get("summary") or {}).get("total_issue_count") or 0,
            storage_root,
        )
    for slide in reversed(category_slides[len(category_details):]):
        _delete_slide(prs, slide)

    key_detail_counts = {
        item.get("name"): int(item.get("count") or 0)
        for item in (report.get("key_issue_summary") or {}).get("details") or []
    }
    empty_key_slides = []
    if not key_detail_counts.get("重点商品"):
        empty_key_slides.extend([prs.slides[24], prs.slides[25], prs.slides[26]])
    if not key_detail_counts.get("月度盘点"):
        empty_key_slides.append(prs.slides[27])
    if not key_detail_counts.get("团购问题"):
        empty_key_slides.append(prs.slides[28])
    if not key_detail_counts.get("商品过期"):
        empty_key_slides.append(prs.slides[29])
    for slide in reversed(empty_key_slides):
        _delete_slide(prs, slide)

    _renumber_slides(prs)

    source_unit_slide = prs.slides[9]
    old_unit_slides = [prs.slides[index] for index in range(9, 20)]
    units = sorted(report.get("units") or [], key=_unit_sort_key)
    new_unit_slides = [
        _add_unit_slide(prs, source_unit_slide, unit)
        for unit in units
    ]
    for slide in reversed(old_unit_slides):
        _delete_slide(prs, slide)
    for index, slide in enumerate(new_unit_slides):
        _move_slide(prs, slide, 9 + index)

    # Defer deletion until template-indexed edits and insertions are complete.
    _delete_slide(prs, detail_slide)
    _renumber_slides(prs)
    _normalize_presentation_fonts(prs)
    _remove_presentation_comments(prs)
    prs.save(output_path)
    slide_files = _render_presentation_preview(output_path, output_dir)
    return {
        "slide_count": len(slide_files),
        "slide_files": slide_files,
        "ppt_path": str(output_path),
    }


def copy_existing_non_oil_presentation(report, destination, storage_root=None):
    presentation = report.get("presentation") or {}
    source = str(presentation.get("ppt_path") or "").strip()
    if source and not os.path.isabs(source) and storage_root:
        source = os.path.join(storage_root, source)
    if not source or not os.path.isfile(source):
        return None
    os.makedirs(os.path.dirname(os.path.abspath(destination)), exist_ok=True)
    # Previously saved reports may still contain template review comments.
    # Strip those on export too, without changing the stored historical deck.
    with ZipFile(source) as package:
        has_comments = any("comment" in name.lower() for name in package.namelist())
    if has_comments:
        deck = Presentation(source)
        _remove_presentation_comments(deck)
        deck.save(destination)
    else:
        shutil.copy2(source, destination)
    return {"slide_count": int(presentation.get("slide_count") or 0)}
