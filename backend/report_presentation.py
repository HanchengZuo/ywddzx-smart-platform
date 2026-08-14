"""Build editable PowerPoint files from saved inspection-report snapshots."""

from __future__ import annotations

import base64
import math
import os
import re
import tempfile
import zipfile
from io import BytesIO
from typing import Sequence

from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
FONT_FAMILY = "Microsoft YaHei"

INK = RGBColor(15, 23, 42)
SLATE = RGBColor(71, 85, 105)
MUTED = RGBColor(100, 116, 139)
LINE = RGBColor(226, 232, 240)
PAPER = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
BLUE = RGBColor(14, 116, 144)
SKY = RGBColor(14, 165, 233)
TEAL = RGBColor(13, 148, 136)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
CHART_COLORS = [
    RGBColor(14, 116, 144),
    RGBColor(20, 184, 166),
    RGBColor(245, 158, 11),
    RGBColor(59, 130, 246),
    RGBColor(239, 68, 68),
    RGBColor(100, 116, 139),
]


def _text(value, fallback="-"):
    normalized = re.sub(r"\s+", " ", str(value or "")).strip()
    return normalized or fallback


def _short(value, limit=120):
    value = _text(value)
    return value if len(value) <= limit else f"{value[: max(1, limit - 1)]}…"


def _chunks(values: Sequence, size: int):
    for index in range(0, len(values), size):
        yield values[index : index + size]


def _split_text(value, limit=280):
    """Split long Chinese narrative text without dropping report content."""
    normalized = _text(value, "")
    if not normalized:
        return []
    sentences = [item.strip() for item in re.split(r"(?<=[。！？；])", normalized) if item.strip()]
    if not sentences:
        sentences = [normalized]
    chunks = []
    current = ""
    for sentence in sentences:
        while len(sentence) > limit:
            if current:
                chunks.append(current)
                current = ""
            chunks.append(sentence[:limit])
            sentence = sentence[limit:]
        if not current:
            current = sentence
        elif len(current) + len(sentence) <= limit:
            current += sentence
        else:
            chunks.append(current)
            current = sentence
    if current:
        chunks.append(current)
    return chunks


def _as_int(value):
    try:
        return int(value or 0)
    except (TypeError, ValueError):
        return 0


def _as_float(value):
    try:
        return float(value or 0)
    except (TypeError, ValueError):
        return 0.0


def _format_metric(value):
    if isinstance(value, float) and not value.is_integer():
        return f"{value:.1f}"
    return str(_as_int(value))


class InspectionReportPresentation:
    def __init__(self, report_type, report, storage_root):
        self.report_type = str(report_type or "").strip()
        self.report = report if isinstance(report, dict) else {}
        self.storage_root = os.path.abspath(storage_root)
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self.page_number = 0

    def build(self, output_path):
        if self.report_type == "quality_measurement":
            self._build_quality_measurement()
            os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
            self.prs.save(output_path)
            self._normalize_chart_axis_ids(output_path)
            return {"slide_count": len(self.prs.slides)}

        self._add_cover()
        self._add_scope_overview()
        builders = {
            "quality_measurement": self._build_quality_measurement,
            "safety_quality": self._build_safety_quality,
            "finance": self._build_finance,
            "on_site_service": self._build_on_site_service,
            "equipment_facilities": self._build_equipment_facilities,
            "non_oil": self._build_non_oil,
        }
        builder = builders.get(self.report_type)
        if not builder:
            raise ValueError("当前报告类型暂不支持导出PPT。")
        builder()
        self._add_ending()
        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        self.prs.save(output_path)
        self._normalize_chart_axis_ids(output_path)
        return {"slide_count": len(self.prs.slides)}

    @staticmethod
    def _normalize_chart_axis_ids(output_path):
        """Convert python-pptx signed axis IDs into spec-compliant UInt32 values."""
        output_path = os.path.abspath(output_path)
        with tempfile.NamedTemporaryFile(
            prefix="report_ppt_",
            suffix=".pptx",
            dir=os.path.dirname(output_path),
            delete=False,
        ) as temp_file:
            temp_path = temp_file.name
        try:
            with zipfile.ZipFile(output_path, "r") as source_zip, zipfile.ZipFile(
                temp_path, "w", zipfile.ZIP_DEFLATED
            ) as target_zip:
                for item in source_zip.infolist():
                    payload = source_zip.read(item.filename)
                    if item.filename.startswith("ppt/charts/chart") and item.filename.endswith(".xml"):
                        payload = re.sub(
                            rb'(axId|crossAx) val="(-\d+)"',
                            lambda match: (
                                match.group(1)
                                + f' val="{int(match.group(2)) + 2 ** 32}"'.encode("ascii")
                            ),
                            payload,
                        )
                    if item.filename.startswith("ppt/slides/slide") and item.filename.endswith(".xml"):
                        payload = re.sub(
                            rb'(<a:latin typeface="([^"]+)"\s*/>)(?!<a:ea)',
                            lambda match: match.group(1) + b'<a:ea typeface="' + match.group(2) + b'"/>',
                            payload,
                        )
                    target_zip.writestr(item, payload)
            os.replace(temp_path, output_path)
        finally:
            if os.path.exists(temp_path):
                os.remove(temp_path)

    def _blank_slide(self, background=PAPER, footer=True):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background_fill = slide.background.fill
        background_fill.solid()
        background_fill.fore_color.rgb = background
        if footer:
            self.page_number += 1
            self._add_footer(slide)
        return slide

    def _add_footer(self, slide):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.55),
            Inches(7.1),
            Inches(12.23),
            Inches(0.012),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = LINE
        shape.line.fill.background()
        self._add_text(
            slide,
            f"业务督导中心  ·  {_text(self.report.get('title'), '巡检分析报告')}",
            0.62,
            7.14,
            10.8,
            0.2,
            size=8,
            color=MUTED,
        )
        self._add_text(
            slide,
            f"{self.page_number:02d}",
            11.75,
            7.11,
            0.9,
            0.25,
            size=9,
            bold=True,
            color=BLUE,
            align=PP_ALIGN.RIGHT,
        )

    def _add_text(
        self,
        slide,
        value,
        x,
        y,
        width,
        height,
        *,
        size=16,
        color=INK,
        bold=False,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.TOP,
        margin=0.04,
    ):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(margin)
        frame.margin_right = Inches(margin)
        frame.margin_top = Inches(margin)
        frame.margin_bottom = Inches(margin)
        frame.vertical_anchor = valign
        paragraph = frame.paragraphs[0]
        paragraph.text = str(value or "")
        paragraph.alignment = align
        paragraph.font.name = FONT_FAMILY
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        return box

    def _add_rich_lines(self, slide, lines, x, y, width, height, size=14, bullet=False):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.08)
        frame.margin_right = Inches(0.08)
        frame.margin_top = Inches(0.06)
        frame.margin_bottom = Inches(0.06)
        for index, line in enumerate(lines):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = _text(line, "")
            paragraph.font.name = FONT_FAMILY
            paragraph.font.size = Pt(size)
            paragraph.font.color.rgb = SLATE
            paragraph.space_after = Pt(7)
            if bullet:
                paragraph.text = f"• {paragraph.text}"
        return box

    def _add_title(self, slide, title, kicker="REPORT", subtitle="", ai=False):
        accent = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.6),
            Inches(0.55),
            Inches(0.11),
            Inches(0.76),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = TEAL
        accent.line.fill.background()
        normalized_title = _text(title, "报告内容")
        title_size = 30 if len(normalized_title) <= 22 else 26
        self._add_text(slide, kicker.upper(), 0.86, 0.5, 5.8, 0.26, size=10, color=TEAL, bold=True)
        self._add_text(slide, normalized_title, 0.83, 0.76, 10.35, 0.64, size=title_size, bold=True)
        if subtitle:
            self._add_text(slide, _short(subtitle, 150), 0.85, 1.4, 11.55, 0.36, size=11, color=MUTED)
        if ai:
            self._add_ai_badge(slide)

    def _add_ai_badge(self, slide, label="AI生成内容"):
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(11.25),
            Inches(0.6),
            Inches(1.35),
            Inches(0.38),
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(224, 242, 254)
        badge.line.color.rgb = RGBColor(125, 211, 252)
        self._set_shape_text(badge, label, 9, BLUE, True)

    def _set_shape_text(self, shape, value, size, color=INK, bold=False, align=PP_ALIGN.CENTER):
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.05)
        frame.margin_right = Inches(0.05)
        frame.margin_top = Inches(0.02)
        frame.margin_bottom = Inches(0.02)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.text = str(value or "")
        paragraph.alignment = align
        paragraph.font.name = FONT_FAMILY
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color

    def _add_cover(self):
        slide = self._blank_slide(background=INK, footer=False)
        glow = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(9.1),
            Inches(0),
            Inches(4.15),
            Inches(4.15),
        )
        glow.fill.solid()
        glow.fill.fore_color.rgb = BLUE
        glow.fill.transparency = 35
        glow.line.fill.background()
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8),
            Inches(0.72),
            Inches(1.2),
            Inches(0.08),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(45, 212, 191)
        line.line.fill.background()
        self._add_text(slide, "AI INSPECTION REPORT", 0.8, 0.92, 5.5, 0.3, size=10, color=RGBColor(94, 234, 212), bold=True)
        self._add_text(
            slide,
            _text(self.report.get("title"), "巡检分析报告"),
            0.78,
            1.45,
            10.8,
            1.35,
            size=34,
            color=WHITE,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )
        self._add_text(
            slide,
            _text(self.report.get("month_label"), "月度报告"),
            0.82,
            2.86,
            6.5,
            0.5,
            size=19,
            color=RGBColor(186, 230, 253),
        )
        scope = self.report.get("source_selection") or {}
        scope_label = "全部可用站点" if scope.get("mode") != "custom" else f"自定义 {len(scope.get('station_ids') or [])} 个站点"
        meta = self.report.get("snapshot") or {}
        generated_at = meta.get("generated_at") or (self.report.get("summary") or {}).get("generated_at") or "-"
        self._add_text(slide, f"数据范围  {scope_label}", 0.82, 5.72, 4.4, 0.36, size=11, color=RGBColor(203, 213, 225))
        self._add_text(slide, f"报告生成  {generated_at}", 0.82, 6.12, 5.5, 0.36, size=11, color=RGBColor(203, 213, 225))
        self._add_text(slide, "业务督导中心数智管理平台", 8.15, 6.18, 4.3, 0.36, size=10, color=RGBColor(148, 163, 184), align=PP_ALIGN.RIGHT)

    def _add_scope_overview(self):
        slide = self._blank_slide()
        self._add_title(slide, "报告口径与数据范围", "DATA SCOPE")
        note = _text(self.report.get("data_scope_note"), "按当前报告模板与账号数据权限统计。")
        self._add_panel(slide, 0.82, 1.9, 11.7, 1.18, "统计口径", note, TEAL)
        source = self.report.get("source_selection") or {}
        metrics = [
            ("统计站点", source.get("station_count") or (self.report.get("summary") or {}).get("station_count"), "座"),
            ("涉及片区", source.get("region_count") or (self.report.get("summary") or {}).get("region_count"), "个"),
            ("巡检记录", source.get("inspection_count"), "条"),
            ("审核通过问题", source.get("issue_count") or (self.report.get("summary") or {}).get("total_issue_count"), "项"),
        ]
        self._add_kpis(slide, metrics, y=3.42)
        targets = self.report.get("target_tables") or []
        self._add_text(slide, "关联检查表", 0.85, 5.22, 2, 0.3, size=11, color=MUTED, bold=True)
        self._add_rich_lines(slide, [f"{index + 1:02d}  {_text(name)}" for index, name in enumerate(targets)], 0.85, 5.58, 11.4, 1.08, size=14)

    def _add_panel(self, slide, x, y, width, height, heading, body, accent=BLUE):
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
        panel.fill.solid()
        panel.fill.fore_color.rgb = WHITE
        panel.line.color.rgb = LINE
        marker = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.22), Inches(y + 0.23), Inches(0.08), Inches(height - 0.46))
        marker.fill.solid()
        marker.fill.fore_color.rgb = accent
        marker.line.fill.background()
        self._add_text(slide, heading, x + 0.46, y + 0.2, width - 0.7, 0.28, size=11, color=accent, bold=True)
        body_text = _text(body)
        body_size = 16 if len(body_text) <= 150 else 14 if len(body_text) <= 260 else 12
        self._add_text(slide, body_text, x + 0.46, y + 0.52, width - 0.7, height - 0.68, size=body_size, color=SLATE, valign=MSO_ANCHOR.MIDDLE)

    def _add_kpis(self, slide, metrics, y=2.0):
        metrics = list(metrics)
        gap = 0.17
        width = (11.7 - gap * max(0, len(metrics) - 1)) / max(1, len(metrics))
        for index, (label, value, unit) in enumerate(metrics):
            x = 0.82 + index * (width + gap)
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(1.38))
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = LINE
            self._add_text(slide, label, x + 0.18, y + 0.18, width - 0.36, 0.3, size=10, color=MUTED, bold=True)
            self._add_text(slide, f"{_format_metric(value)}{unit}", x + 0.16, y + 0.52, width - 0.32, 0.62, size=25, color=BLUE, bold=True, valign=MSO_ANCHOR.MIDDLE)

    def _add_section(self, chapter, title, subtitle=""):
        slide = self._blank_slide(background=INK)
        self._add_text(slide, chapter, 0.86, 1.3, 2.7, 0.34, size=11, color=RGBColor(94, 234, 212), bold=True)
        self._add_text(slide, title, 0.83, 1.9, 10.8, 1.05, size=31, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
        if subtitle:
            self._add_text(slide, subtitle, 0.85, 3.15, 10.6, 1.2, size=16, color=RGBColor(203, 213, 225))
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.86), Inches(5.95), Inches(4.2), Inches(0.08))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = TEAL
        stripe.line.fill.background()
        return slide

    def _add_narrative_slide(self, title, narrative, *, kicker="OVERVIEW", ai=False, metrics=None):
        slide = self._blank_slide()
        self._add_title(slide, title, kicker, ai=ai)
        self._add_panel(slide, 0.82, 1.9, 11.7, 1.78, "情况概述", _text(narrative), TEAL)
        if metrics:
            self._add_kpis(slide, metrics, y=4.02)
        return slide

    def _add_insight_strip(self, slide, narrative, y=1.56):
        body = _short(narrative, 165)
        strip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.84),
            Inches(y),
            Inches(11.64),
            Inches(0.56),
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = RGBColor(236, 254, 255)
        strip.line.color.rgb = RGBColor(165, 243, 252)
        self._add_text(slide, "核心发现", 1.02, y + 0.13, 0.78, 0.28, size=10, color=BLUE, bold=True)
        self._add_text(slide, body, 1.86, y + 0.11, 10.36, 0.32, size=12, color=SLATE)

    def _add_chart_slides(self, title, items, *, name_key="name", value_keys=("count",), series_names=("问题数量",), kicker="DATA ANALYSIS", narrative="", chunk_size=12):
        items = [item for item in (items or []) if isinstance(item, dict)]
        if not items:
            self._add_empty_slide(title, "当前范围暂无可展示的统计数据。", kicker)
            return
        longest_label = max(len(_text(item.get(name_key), "")) for item in items)
        effective_chunk_size = min(chunk_size, 8 if longest_label >= 9 else 10)
        for page_index, chunk in enumerate(_chunks(items, effective_chunk_size), 1):
            slide = self._blank_slide()
            page_title = title if len(items) <= effective_chunk_size else f"{title}（{page_index}）"
            self._add_title(slide, page_title, kicker)
            page_narrative = narrative if page_index == 1 else ""
            if page_narrative:
                self._add_insight_strip(slide, page_narrative)
            chart_data = ChartData()
            chart_data.categories = [_short(item.get(name_key), 18) for item in chunk]
            all_values = []
            for series_index, value_key in enumerate(value_keys):
                values = [_as_float(item.get(value_key)) for item in chunk]
                all_values.extend(values)
                chart_data.add_series(series_names[series_index], values)
            use_bar_chart = longest_label >= 8 or len(chunk) >= 9
            chart_type = XL_CHART_TYPE.BAR_CLUSTERED if use_bar_chart else XL_CHART_TYPE.COLUMN_CLUSTERED
            y = 2.3 if page_narrative else 1.78
            height = 4.32 if page_narrative else 4.86
            chart = slide.shapes.add_chart(
                chart_type,
                Inches(0.86),
                Inches(y),
                Inches(11.55),
                Inches(height),
                chart_data,
            ).chart
            chart.has_legend = len(value_keys) > 1
            if chart.has_legend:
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
                chart.legend.font.name = FONT_FAMILY
                chart.legend.font.size = Pt(10)
            chart.value_axis.minimum_scale = 0
            chart.value_axis.has_major_gridlines = True
            chart.value_axis.major_gridlines.format.line.color.rgb = LINE
            chart.value_axis.format.line.color.rgb = RGBColor(203, 213, 225)
            chart.value_axis.tick_labels.font.name = FONT_FAMILY
            chart.value_axis.tick_labels.font.size = Pt(10)
            chart.category_axis.format.line.color.rgb = RGBColor(203, 213, 225)
            chart.category_axis.tick_labels.font.name = FONT_FAMILY
            chart.category_axis.tick_labels.font.size = Pt(10)
            if use_bar_chart:
                chart.category_axis.reverse_order = True
            has_decimal = any(value and not float(value).is_integer() for value in all_values)
            for series_index, series in enumerate(chart.series):
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = CHART_COLORS[series_index % len(CHART_COLORS)]
                series.format.line.fill.background()
            for plot in chart.plots:
                plot.gap_width = 72 if use_bar_chart else 58
                plot.has_data_labels = True
                labels = plot.data_labels
                labels.position = XL_LABEL_POSITION.OUTSIDE_END
                labels.font.name = FONT_FAMILY
                labels.font.size = Pt(9 if len(value_keys) > 1 else 10)
                labels.font.bold = True
                labels.font.color.rgb = SLATE
                labels.number_format = "0.0" if has_decimal else "0"
                labels.number_format_is_linked = False

    @staticmethod
    def _is_narrative_column(header):
        return any(keyword in _text(header, "") for keyword in ("描述", "摘要", "规定", "内容", "说明", "分析", "建议"))

    def _table_column_weights(self, headers, rows):
        weights = []
        for index, header in enumerate(headers):
            name = _text(header, "")
            values = [_text(row[index] if index < len(row) else "", "") for row in rows]
            average_length = sum(len(value) for value in values) / max(1, len(values))
            if self._is_narrative_column(name):
                weight = 3.6
            elif any(keyword in name for keyword in ("问题数量", "问题合计", "受检站点", "排名", "类型")):
                weight = 0.9
            elif any(keyword in name for keyword in ("站点", "单位", "项目", "环节", "日期", "时间")):
                weight = 1.45
            else:
                weight = min(2.2, max(1.0, 0.9 + average_length / 18))
            weights.append(weight)
        return weights

    def _table_row_units(self, row, headers, weights):
        required_lines = 1
        for index, value in enumerate(row):
            text_value = _text(value, "")
            characters_per_line = max(7, int(13 * weights[index]))
            required_lines = max(required_lines, math.ceil(len(text_value) / characters_per_line))
        return max(1, min(4, required_lines))

    def _paginate_table_rows(self, rows, headers, weights, unit_limit):
        pages = []
        current = []
        current_units = 0
        for row in rows:
            units = self._table_row_units(row, headers, weights)
            if current and current_units + units > unit_limit:
                pages.append(current)
                current = []
                current_units = 0
            current.append((row, units))
            current_units += units
        if current:
            pages.append(current)
        return pages

    def _add_table_slides(self, title, headers, rows, *, kicker="DATA TABLE", rows_per_slide=11, ai=False):
        normalized_rows = [list(row) for row in (rows or [])]
        if not normalized_rows:
            self._add_empty_slide(title, "当前范围暂无明细数据。", kicker)
            return
        weights = self._table_column_weights(headers, normalized_rows)
        pages = self._paginate_table_rows(normalized_rows, headers, weights, rows_per_slide)
        total_width = int(Inches(11.88))
        allocated_width = 0
        for page_index, page in enumerate(pages, 1):
            chunk = [row for row, _ in page]
            row_units = [units for _, units in page]
            slide = self._blank_slide()
            page_title = title if len(pages) == 1 else f"{title}（{page_index}/{len(pages)}）"
            self._add_title(slide, page_title, kicker, ai=ai)
            shape = slide.shapes.add_table(
                len(chunk) + 1,
                len(headers),
                Inches(0.72),
                Inches(1.78),
                Inches(11.88),
                Inches(4.98),
            )
            table = shape.table
            total_weight = sum(weights)
            allocated_width = 0
            for column_index, column in enumerate(table.columns):
                if column_index == len(table.columns) - 1:
                    column_width = total_width - allocated_width
                else:
                    column_width = int(total_width * weights[column_index] / total_weight)
                    allocated_width += column_width
                column.width = column_width
            table.rows[0].height = Inches(0.54)
            available_height = 4.44
            unit_height = available_height / max(1, sum(row_units))
            for row_index, units in enumerate(row_units, 1):
                table.rows[row_index].height = Inches(unit_height * units)
            for column_index, header in enumerate(headers):
                cell = table.cell(0, column_index)
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
                self._set_cell_text(cell, header, 10, WHITE, True)
            for row_index, row in enumerate(chunk, 1):
                for column_index in range(len(headers)):
                    cell = table.cell(row_index, column_index)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE if row_index % 2 else RGBColor(241, 245, 249)
                    value = row[column_index] if column_index < len(row) else ""
                    narrative_column = self._is_narrative_column(headers[column_index])
                    body_size = 9 if len(_text(value, "")) > 110 else 10
                    self._set_cell_text(
                        cell,
                        value,
                        body_size,
                        SLATE,
                        align=PP_ALIGN.LEFT if narrative_column else PP_ALIGN.CENTER,
                    )

    def _set_cell_text(self, cell, value, size, color, bold=False, align=PP_ALIGN.CENTER):
        cell.margin_left = Inches(0.05)
        cell.margin_right = Inches(0.05)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        frame = cell.text_frame
        frame.clear()
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.text = str(value or "")
        paragraph.alignment = align
        paragraph.font.name = FONT_FAMILY
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color

    def _add_empty_slide(self, title, message, kicker="REPORT"):
        slide = self._blank_slide()
        self._add_title(slide, title, kicker)
        self._add_panel(slide, 2.15, 2.65, 9.0, 1.55, "暂无数据", message, MUTED)

    def _resolve_image_source(self, raw_path):
        value = str(raw_path or "").strip()
        if not value:
            return None
        if value.startswith("data:image/") and ";base64," in value:
            try:
                return BytesIO(base64.b64decode(value.split(",", 1)[1], validate=True))
            except (ValueError, TypeError):
                return None
        clean = value.split("?", 1)[0].replace("\\", "/")
        if "/storage/" in clean:
            clean = clean.split("/storage/", 1)[1]
        clean = clean.lstrip("/")
        candidate = os.path.abspath(os.path.join(self.storage_root, clean))
        try:
            if os.path.commonpath([self.storage_root, candidate]) != self.storage_root:
                return None
        except ValueError:
            return None
        return candidate if os.path.isfile(candidate) else None

    def _add_picture_contain(self, slide, image_source, x, y, width, height):
        frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(width),
            Inches(height),
        )
        frame.fill.solid()
        frame.fill.fore_color.rgb = RGBColor(241, 245, 249)
        frame.line.color.rgb = RGBColor(203, 213, 225)
        source = self._resolve_image_source(image_source)
        if not source:
            self._set_shape_text(frame, "暂无问题照片", 11, MUTED)
            return
        try:
            if hasattr(source, "seek"):
                source.seek(0)
            with Image.open(source) as image:
                image_width, image_height = image.size
            if hasattr(source, "seek"):
                source.seek(0)
            padding = 0.07
            inner_width = max(0.1, width - padding * 2)
            inner_height = max(0.1, height - padding * 2)
            ratio = min(inner_width / image_width, inner_height / image_height)
            target_width = image_width * ratio
            target_height = image_height * ratio
            left = x + (width - target_width) / 2
            top = y + (height - target_height) / 2
            slide.shapes.add_picture(source, Inches(left), Inches(top), Inches(target_width), Inches(target_height))
        except Exception:
            frame.fill.fore_color.rgb = RGBColor(254, 242, 242)
            frame.line.color.rgb = RGBColor(254, 202, 202)
            self._set_shape_text(frame, "照片读取失败", 11, RED)

    @staticmethod
    def _issue_description(issue):
        return _text(
            issue.get("description")
            or issue.get("inspection_content")
            or issue.get("summary"),
            "暂无问题描述",
        )

    def _collect_nested_issues(self, value, *, defaults=None):
        """Collect selected issue records from report-specific nested analysis data."""
        defaults = dict(defaults or {})
        collected = []
        if isinstance(value, list):
            for item in value:
                collected.extend(self._collect_nested_issues(item, defaults=defaults))
            return collected
        if not isinstance(value, dict):
            return collected

        inherited = dict(defaults)
        for source_key, target_key in (
            ("unit_name", "unit_name"),
            ("station_name", "station_name"),
            ("area_name", "service_area"),
            ("service_area", "service_area"),
            ("category_name", "category_name"),
            ("title", "category_name"),
        ):
            if value.get(source_key) and not inherited.get(target_key):
                inherited[target_key] = value.get(source_key)

        is_issue = bool(
            value.get("issue_id")
            or value.get("issue_photo")
            or value.get("description")
            or value.get("inspection_content")
        )
        if is_issue:
            issue = dict(value)
            for key, default_value in inherited.items():
                issue.setdefault(key, default_value)
            collected.append(issue)
            return collected

        for child_key in (
            "issues",
            "highlighted_issues",
            "highlights",
            "areas",
            "area_analyses",
            "service_areas",
        ):
            child_value = value.get(child_key)
            if child_value:
                collected.extend(self._collect_nested_issues(child_value, defaults=inherited))
        return collected

    def _issue_pages(self, issues):
        pages = []
        paired = []
        for issue in issues:
            if len(self._issue_description(issue)) > 190:
                if paired:
                    pages.append(paired)
                    paired = []
                pages.append([issue])
                continue
            paired.append(issue)
            if len(paired) == 2:
                pages.append(paired)
                paired = []
        if paired:
            pages.append(paired)
        return pages

    def _add_issue_slides(self, title, issues, *, kicker="TYPICAL ISSUES", ai=False, subtitle="", max_issues=12):
        issues = [item for item in (issues or []) if isinstance(item, dict)][:max_issues]
        if not issues:
            self._add_empty_slide(title, "当前范围暂无可展示的问题。", kicker)
            return
        pages = self._issue_pages(issues)
        for page_index, chunk in enumerate(pages, 1):
            slide = self._blank_slide()
            page_title = title if len(pages) == 1 else f"{title}（{page_index}/{len(pages)}）"
            self._add_title(slide, page_title, kicker, subtitle=subtitle if page_index == 1 else "", ai=ai)
            start_y = 1.98 if subtitle and page_index == 1 else 1.78
            if len(chunk) == 1:
                issue = chunk[0]
                station = issue.get("station_name") or issue.get("unit_name") or issue.get("management_unit") or "问题明细"
                category = issue.get("category_name") or issue.get("inspection_item") or issue.get("project") or issue.get("service_area") or ""
                self._add_picture_contain(slide, issue.get("issue_photo"), 0.82, start_y, 5.45, 4.86)
                self._add_text(slide, _short(station, 28), 6.6, start_y + 0.1, 5.65, 0.5, size=20, color=BLUE, bold=True)
                if category:
                    self._add_text(slide, _short(category, 42), 6.62, start_y + 0.68, 5.55, 0.36, size=13, color=TEAL, bold=True)
                self._add_text(slide, "原始问题描述", 6.62, start_y + 1.24, 2.1, 0.3, size=11, color=MUTED, bold=True)
                description = self._issue_description(issue)
                description_size = 16 if len(description) <= 150 else 14 if len(description) <= 260 else 12
                self._add_text(slide, description, 6.6, start_y + 1.62, 5.68, 2.98, size=description_size, color=SLATE)
                continue

            for index, issue in enumerate(chunk):
                x = 0.72 + index * 6.12
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(x),
                    Inches(start_y),
                    Inches(5.82),
                    Inches(4.9),
                )
                card.fill.solid()
                card.fill.fore_color.rgb = WHITE
                card.line.color.rgb = LINE
                self._add_picture_contain(slide, issue.get("issue_photo"), x + 0.18, start_y + 0.18, 5.46, 2.58)
                station = issue.get("station_name") or issue.get("unit_name") or issue.get("management_unit") or "问题明细"
                category = issue.get("category_name") or issue.get("inspection_item") or issue.get("project") or issue.get("service_area") or ""
                self._add_text(slide, _short(station, 24), x + 0.24, start_y + 2.94, 5.34, 0.38, size=15, color=BLUE, bold=True)
                if category:
                    self._add_text(slide, _short(category, 38), x + 0.24, start_y + 3.36, 5.34, 0.3, size=11, color=TEAL, bold=True)
                description = self._issue_description(issue)
                self._add_text(slide, _short(description, 190), x + 0.23, start_y + 3.72, 5.36, 0.95, size=11, color=SLATE)

    def _add_analysis_cards(self, title, items, *, kicker="AI ANALYSIS", ai=True):
        items = [item for item in (items or []) if isinstance(item, dict)]
        if not items:
            self._add_empty_slide(title, "当前报告暂无可展示的分析内容。", kicker)
            return
        expanded_items = []
        for item in items:
            item_title = _text(item.get("title"), "分析事项")
            content_chunks = _split_text(item.get("content"), 270) or ["暂无详细内容"]
            for chunk_index, content in enumerate(content_chunks):
                expanded_items.append(
                    {
                        "title": item_title if chunk_index == 0 else f"{item_title}（续）",
                        "content": content,
                        "ai_generated": item.get("ai_generated", True),
                    }
                )
        pages = list(_chunks(expanded_items, 2))
        for page_index, chunk in enumerate(pages, 1):
            slide = self._blank_slide()
            page_title = title if len(pages) == 1 else f"{title}（{page_index}/{len(pages)}）"
            self._add_title(slide, page_title, kicker, ai=ai and any(item.get("ai_generated", True) for item in chunk))
            gap = 0.2
            card_height = (4.94 - gap * (len(chunk) - 1)) / len(chunk)
            for index, item in enumerate(chunk):
                y = 1.78 + index * (card_height + gap)
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.82),
                    Inches(y),
                    Inches(11.7),
                    Inches(card_height),
                )
                card.fill.solid()
                card.fill.fore_color.rgb = WHITE
                card.line.color.rgb = LINE
                number = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.02), Inches(y + 0.22), Inches(0.62), Inches(0.62))
                number.fill.solid()
                number.fill.fore_color.rgb = RGBColor(204, 251, 241)
                number.line.fill.background()
                self._set_shape_text(number, f"{(page_index - 1) * 2 + index + 1:02d}", 12, TEAL, True)
                self._add_text(slide, item.get("title"), 1.84, y + 0.18, 10.2, 0.44, size=17 if len(chunk) == 1 else 15, bold=True)
                content = _text(item.get("content"), "暂无详细内容")
                content_size = 16 if len(chunk) == 1 and len(content) <= 180 else 14 if len(content) <= 210 else 13
                self._add_text(slide, content, 1.84, y + 0.72, 10.25, card_height - 0.92, size=content_size, color=SLATE)

    def _add_quality_header(self, slide, title, ai=False):
        self._add_text(slide, title, 0.62, 0.28, 10.9, 0.5, size=25, bold=True)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.92), Inches(13.333), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(42, 155, 211)
        line.line.fill.background()
        logo_path = os.path.join(os.path.dirname(__file__), "assets", "report_logo.png")
        if os.path.isfile(logo_path):
            slide.shapes.add_picture(logo_path, Inches(12.42), Inches(0.18), height=Inches(0.56))
        if ai:
            self._add_ai_badge(slide)

    def _quality_blank_slide(self, title, ai=False):
        slide = self._blank_slide(background=WHITE, footer=False)
        self.page_number += 1
        self._add_quality_header(slide, title, ai=ai)
        self._add_text(slide, str(self.page_number), 12.45, 7.15, 0.48, 0.2, size=8, color=MUTED, align=PP_ALIGN.RIGHT)
        return slide

    def _quality_table(self, slide, headers, rows, x, y, width, height, weights=None, narrative_columns=None):
        rows = [list(row) for row in rows]
        shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(width), Inches(height))
        table = shape.table
        weights = weights or [1] * len(headers)
        total_weight = sum(weights)
        total_width = int(Inches(width))
        allocated = 0
        for index, column in enumerate(table.columns):
            if index == len(headers) - 1:
                column.width = total_width - allocated
            else:
                column.width = int(total_width * weights[index] / total_weight)
                allocated += column.width
        table.rows[0].height = Inches(0.46)
        body_height = max(0.35, (height - 0.46) / max(1, len(rows)))
        for index in range(1, len(table.rows)):
            table.rows[index].height = Inches(body_height)
        narrative_columns = set(narrative_columns or [])
        for column_index, header in enumerate(headers):
            cell = table.cell(0, column_index)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(19, 82, 139)
            self._set_cell_text(cell, header, 9, WHITE, True)
        for row_index, row in enumerate(rows, 1):
            for column_index in range(len(headers)):
                cell = table.cell(row_index, column_index)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if row_index % 2 else RGBColor(242, 247, 251)
                value = row[column_index] if column_index < len(row) else ""
                font_size = 7 if len(_text(value, "")) > 80 else 8
                self._set_cell_text(
                    cell,
                    value,
                    font_size,
                    INK,
                    bold=row_index == len(rows),
                    align=PP_ALIGN.LEFT if column_index in narrative_columns else PP_ALIGN.CENTER,
                )
        return table

    def _render_quality_overall(self, data):
        slide = self._quality_blank_slide(data.get("title") or "总体情况")
        self._add_text(slide, data.get("narrative"), 0.7, 1.12, 11.95, 1.28, size=12, color=INK)
        rows = list(data.get("rows") or []) + [data.get("total_row") or {}]
        table_rows = [[
            item.get("sequence"), item.get("unit_name"), item.get("oil_depot_count", 0),
            item.get("station_count", 0), item.get("transport_vehicle_count", 0),
            item.get("general_issue_count", 0), item.get("violation_issue_count", 0),
            item.get("prohibited_issue_count", 0), item.get("total_issue_count", 0),
        ] for item in rows]
        self._quality_table(
            slide,
            ["序号", "二级单位", "检查油库\n数量", "检查加油站\n数量", "检查运输车辆\n数量", "一般性\n问题", "违规违纪\n问题", "涉及禁止项\n问题", "单库、车、站\n问题数量"],
            table_rows, 0.48, 2.52, 12.36, 4.34,
            weights=[0.55, 1.45, 0.85, 0.9, 0.95, 0.78, 0.78, 0.9, 1.0],
        )

    def _render_quality_finding(self, data):
        slide = self._quality_blank_slide(data.get("title") or "检查发现-发现问题")
        self._add_rich_lines(slide, data.get("text_lines") or [], 0.62, 1.26, 5.08, 5.75, size=12)
        rows = [[item.get("sequence"), item.get("section"), item.get("problem_type"), item.get("count"), item.get("percentage")] for item in data.get("rows") or []]
        self._quality_table(slide, ["序号", "环节排前三", "问题类型", "问题数量", "占比/%"], rows, 5.88, 1.25, 6.84, 5.62, weights=[0.55, 1.15, 1.55, 0.82, 0.75])

    def _render_quality_prohibited(self, data):
        slide = self._quality_blank_slide(data.get("title") or "检查发现-禁止项问题")
        self._add_text(slide, data.get("narrative"), 0.72, 1.15, 11.9, 0.72, size=14, bold=True)
        rows = [[index, "加油站环节", item.get("unit_name"), item.get("description"), item.get("penalty") or ""] for index, item in enumerate(data.get("rows") or [], 1)]
        if not rows:
            rows = [["-", "加油站环节", "-", "当前月份暂无禁止项问题", ""]]
        self._quality_table(slide, ["序号", "环节", "基层单位名称", "禁止项管理规定", "处罚情况"], rows, 0.52, 2.0, 12.28, 4.85, weights=[0.45, 0.9, 1.15, 4.3, 0.9], narrative_columns={3})

    def _render_quality_flow_chart(self, data):
        slide = self._quality_blank_slide(data.get("title") or "检查发现-加油站环节")
        self._add_text(slide, data.get("narrative"), 0.78, 1.14, 11.8, 1.05, size=14, bold=True)
        self._add_text(slide, data.get("chart_title") or "各类问题数量汇总情况", 3.8, 2.18, 5.8, 0.45, size=19, bold=True, align=PP_ALIGN.CENTER)
        distribution = data.get("distribution") or []
        if not distribution:
            self._add_panel(slide, 2.1, 3.25, 9.1, 1.5, "暂无数据", "当前月份暂无业务流程分布数据。", MUTED)
            return
        chart_data = ChartData()
        chart_data.categories = [_short(item.get("name"), 14) for item in distribution]
        chart_data.add_series("问题数量", [_as_int(item.get("count")) for item in distribution])
        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.05), Inches(2.75), Inches(11.15), Inches(3.75), chart_data).chart
        chart.has_legend = False
        chart.value_axis.minimum_scale = 0
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = LINE
        chart.value_axis.tick_labels.font.name = FONT_FAMILY
        chart.value_axis.tick_labels.font.size = Pt(10)
        chart.category_axis.tick_labels.font.name = FONT_FAMILY
        chart.category_axis.tick_labels.font.size = Pt(11)
        series = chart.series[0]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = RGBColor(76, 157, 213)
        series.format.line.fill.background()
        plot = chart.plots[0]
        plot.gap_width = 75
        plot.has_data_labels = True
        plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        plot.data_labels.font.name = FONT_FAMILY
        plot.data_labels.font.size = Pt(11)
        plot.data_labels.font.bold = True

    def _render_quality_issue_pairs(self, data):
        title = data.get("title") or "加油站环节"
        if int(data.get("continuation_count") or 1) > 1:
            title = f"{title}（{data.get('continuation')}/{data.get('continuation_count')}）"
        slide = self._quality_blank_slide(title)
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.34), Inches(1.1), Inches(12.45), Inches(0.54))
        band.fill.solid()
        band.fill.fore_color.rgb = RGBColor(42, 137, 193)
        band.line.fill.background()
        self._add_text(slide, data.get("subtitle"), 0.48, 1.19, 11.9, 0.3, size=13, color=WHITE, bold=True)
        issues = data.get("issues") or []
        for index, issue in enumerate(issues):
            x = 0.58 + index * 6.18
            if index:
                divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.63), Inches(1.9), Inches(0.01), Inches(4.75))
                divider.fill.solid(); divider.fill.fore_color.rgb = RGBColor(151, 203, 230); divider.line.fill.background()
            self._add_text(slide, _short(issue.get("station_name"), 28), x, 1.9, 5.65, 0.42, size=16, bold=True)
            self._add_text(slide, issue.get("description"), x, 2.38, 5.65, 1.25, size=12, color=INK)
            self._add_picture_contain(slide, issue.get("issue_photo"), x + 1.15, 3.82, 3.35, 2.65)
        if not issues:
            self._add_panel(slide, 2.1, 3.0, 9.0, 1.5, "暂无数据", "当前分类暂无可展示的问题。", MUTED)

    def _render_quality_management_trace(self, data):
        slide = self._quality_blank_slide(data.get("title") or "管理追溯", ai=bool(data.get("ai_generated")))
        typical_issue = data.get("typical_issue") or {}
        description = self._issue_description(typical_issue)
        station = _text(typical_issue.get("station_name"), "典型问题")
        self._add_text(slide, "典型问题：", 0.72, 1.25, 1.55, 0.36, size=15, color=RED, bold=True)
        self._add_text(slide, f"{station}：{description}", 1.85, 1.25, 7.2, 0.82, size=14, color=RED, bold=True)
        y = 2.18
        for item in data.get("analysis_items") or []:
            self._add_text(slide, item.get("label"), 0.88, y, 1.62, 0.36, size=13, color=RGBColor(52, 119, 195), bold=True)
            self._add_text(slide, item.get("content"), 2.22, y, 6.75, 0.95, size=12, color=INK)
            y += 1.08
        self._add_picture_contain(slide, typical_issue.get("issue_photo"), 9.28, 1.58, 3.25, 4.92)

    def _render_quality_trace_analysis(self, data):
        slide = self._quality_blank_slide(data.get("title") or "管理追溯", ai=bool(data.get("ai_generated")))
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.12), Inches(1.12), Inches(8.2), Inches(0.6))
        band.fill.solid(); band.fill.fore_color.rgb = RGBColor(42, 137, 193); band.line.fill.background()
        self._add_text(slide, data.get("subtitle") or "典型问题分析", 0.28, 1.23, 7.6, 0.32, size=16, color=WHITE, bold=True)
        self._add_text(slide, "综上所述：", 0.92, 2.06, 1.5, 0.36, size=14, color=RGBColor(52, 119, 195), bold=True)
        conclusion = _text(data.get("conclusion"), "暂无分析结论。")
        if conclusion.startswith("综上所述："):
            conclusion = conclusion[len("综上所述："):]
        self._add_text(slide, conclusion, 2.18, 2.04, 10.1, 1.2, size=14, color=INK)
        self._add_text(slide, "改进措施：", 0.92, 3.42, 1.5, 0.36, size=14, color=RGBColor(52, 119, 195), bold=True)
        lines = [f"{index}、{item.get('level')}：{item.get('content')}" for index, item in enumerate(data.get("improvement_measures") or [], 1)]
        self._add_rich_lines(slide, lines or ["暂无改进措施。"], 1.0, 3.85, 11.35, 2.55, size=13)

    def _render_quality_work_plan(self, data):
        slide = self._quality_blank_slide(data.get("title") or "工作计划", ai=bool(data.get("ai_generated")))
        y = 1.3
        for index, item in enumerate(data.get("items") or [], 1):
            self._add_text(slide, f"{index}、{item.get('title')}", 1.05, y, 10.95, 0.36, size=15, color=RGBColor(52, 119, 195), bold=True)
            self._add_text(slide, item.get("content"), 1.05, y + 0.38, 11.1, 0.78, size=12, color=INK)
            y += 1.3

    def _build_quality_measurement(self):
        renderers = {
            "overall": self._render_quality_overall,
            "finding_overview": self._render_quality_finding,
            "prohibited": self._render_quality_prohibited,
            "flow_chart": self._render_quality_flow_chart,
            "issue_pairs": self._render_quality_issue_pairs,
            "management_trace": self._render_quality_management_trace,
            "trace_analysis": self._render_quality_trace_analysis,
            "work_plan": self._render_quality_work_plan,
        }
        slides = [item for item in (self.report.get("slides") or []) if isinstance(item, dict)]
        if not slides:
            self._add_empty_slide("质量计量监督检查报告", "当前报告快照尚未生成幻灯片数据，请重新生成报告。")
            return
        for slide_data in slides:
            renderer = renderers.get(slide_data.get("kind"))
            if renderer:
                renderer(slide_data)

    def _build_safety_quality(self):
        summary = self.report.get("summary") or {}
        self._add_section("CHAPTER 01", "总体情况", "视频扫站与四不两直现场检查分线统计。")
        self._add_narrative_slide(
            "安全质量检查概览",
            "视频扫站和现场检查分别统计，图表中的站点数与审核通过问题数来自当前报告快照。",
            metrics=[
                ("视频站点", summary.get("video_station_count"), "座"),
                ("视频问题", summary.get("video_issue_count"), "项"),
                ("现场站点", summary.get("onsite_station_count"), "座"),
                ("现场问题", summary.get("onsite_issue_count"), "项"),
            ],
        )
        for section in self.report.get("sections") or []:
            self._add_chart_slides(
                f"{section.get('label')} · 单位分布",
                section.get("units") or [],
                name_key="unit_name",
                value_keys=("issue_count", "station_count"),
                series_names=("问题数量", "检查站点"),
                narrative=section.get("narrative") or "",
            )
            self._add_chart_slides(
                f"{section.get('label')} · {section.get('category_field') or '问题分类'}",
                section.get("category_distribution") or [],
                narrative=section.get("category_text") or "",
            )
        deep = self.report.get("deep_analysis") or {}
        self._add_section("CHAPTER 02", "典型问题与重点问题", "AI辅助筛选高频典型问题与分类重点问题。")
        for finding in deep.get("typical_findings") or []:
            self._add_issue_slides(
                f"{finding.get('label') or finding.get('mode') or '典型问题'}",
                finding.get("issues") or ([finding.get("representative_issue")] if finding.get("representative_issue") else []),
                ai=bool(finding.get("ai_generated")),
                subtitle=finding.get("summary") or "",
                max_issues=8,
            )
        for highlight in (deep.get("category_highlights") or [])[:10]:
            self._add_issue_slides(
                f"{highlight.get('label') or ''} · {highlight.get('category_name') or '重点问题'}",
                highlight.get("issues") or [],
                ai=bool(highlight.get("ai_generated")),
                subtitle=highlight.get("summary") or "",
                max_issues=6,
            )
        self._add_section("CHAPTER 03", "问题分析与工作建议")
        self._add_analysis_cards("问题分析", deep.get("problem_analysis") or [], ai=True)
        self._add_analysis_cards("工作建议", deep.get("work_suggestions") or [], ai=True)

    def _build_finance(self):
        summary = self.report.get("summary") or {}
        self._add_section("CHAPTER 01", "总体情况", "按单位、站点、项目和关键环节统计财务检查结果。")
        self._add_narrative_slide(
            "财务检查概览",
            self.report.get("overview_text") or self.report.get("scope_text"),
            metrics=[
                ("检查站点", summary.get("station_count"), "座"),
                ("管理片区", summary.get("region_count"), "个"),
                ("控参股单位", summary.get("holding_unit_count"), "个"),
                ("发现问题", summary.get("total_issue_count"), "项"),
            ],
        )
        units = self.report.get("units") or []
        self._add_chart_slides("单位问题与站点分布", units, name_key="unit_name", value_keys=("issue_count", "station_count"), series_names=("问题数量", "检查站点"))
        self._add_chart_slides("检查项目分布", self.report.get("project_distribution") or [], narrative=self.report.get("project_distribution_text") or "")
        self._add_chart_slides("关键环节分布", self.report.get("key_link_distribution") or [], narrative=self.report.get("key_link_distribution_text") or "")
        self._add_section("CHAPTER 02", "各站结果通报")
        station_rows = self.report.get("station_reports") or []
        self._add_table_slides(
            "站点问题概览",
            ["所属单位", "站点", "巡检时间", "问题数量", "问题摘要"],
            [[item.get("unit_name"), item.get("station_name"), item.get("date_range"), item.get("issue_count"), "；".join(_short(issue.get("description"), 42) for issue in (item.get("issues") or [])[:3])] for item in station_rows],
            rows_per_slide=8,
        )
        photographed = []
        for station in station_rows:
            for issue in station.get("issues") or []:
                issue = dict(issue)
                issue.setdefault("station_name", station.get("station_name"))
                photographed.append(issue)
        self._add_issue_slides("各站代表问题", photographed, max_issues=12)
        deep = self.report.get("deep_analysis") or {}
        self._add_section("CHAPTER 03", "检查结果分析与建议")
        self._add_analysis_cards("检查结果分析", deep.get("result_analysis") or [], ai=True)
        self._add_analysis_cards("检查内容建议", deep.get("content_suggestions") or [], ai=True)

    def _build_equipment_facilities(self):
        summary = self.report.get("summary") or {}
        self._add_section("CHAPTER 01", "总体情况", "按片区与站点两个维度分析设备设施问题。")
        self._add_narrative_slide(
            "设备设施检查概览",
            self.report.get("overview_text"),
            metrics=[
                ("检查单位", summary.get("unit_count"), "个"),
                ("受检站点", summary.get("station_count"), "座"),
                ("发现问题", summary.get("total_issue_count"), "项"),
                ("站均问题", summary.get("average_issue_count"), "项"),
            ],
        )
        self._add_chart_slides("按片区划分", self.report.get("region_rows") or [], name_key="unit_name", value_keys=("issue_count", "station_count", "average_issue_count"), series_names=("问题数量", "受检站点", "平均问题"), narrative=self.report.get("region_text") or "")
        rankings = self.report.get("station_ranking") or []
        self._add_table_slides("站点问题数量排名", ["排名", "站点", "所属单位", "检查日期", "问题数量"], [[item.get("rank"), item.get("station_name"), item.get("management_unit"), item.get("date_range") or item.get("report_date"), item.get("issue_count")] for item in rankings], rows_per_slide=12)
        self._add_section("CHAPTER 02", "问题数据统计分析")
        self._add_chart_slides("所属区域问题分布", self.report.get("area_distribution") or [], narrative=self.report.get("area_distribution_text") or "")
        self._add_chart_slides("检查事项问题分布", self.report.get("item_distribution") or [], narrative=self.report.get("item_distribution_text") or "")
        deep = self.report.get("deep_analysis") or {}
        self._add_section("CHAPTER 03", "典型问题、分析与建议")
        typical = deep.get("typical_finding") or {}
        representative = typical.get("representative_issue")
        self._add_issue_slides("检查发现 · 典型问题", [representative] if representative else [], ai=bool(typical.get("ai_generated")), subtitle=typical.get("summary") or "")
        self._add_analysis_cards("问题分析", deep.get("problem_analysis") or [], ai=True)
        self._add_analysis_cards("工作建议", deep.get("work_suggestions") or [], ai=True)

    def _build_on_site_service(self):
        summary = self.report.get("summary") or {}
        self._add_section("CHAPTER 01", "检查基本情况", "线上与线下巡检分线统计，并结合上月整改情况进行分析。")
        self._add_narrative_slide(
            "现场服务检查概览",
            self.report.get("overview_text"),
            metrics=[
                ("检查站次", summary.get("station_visit_count"), "座"),
                ("发现问题", summary.get("total_issue_count"), "项"),
                ("视频站均", summary.get("video_average_issue_count"), "项"),
                ("现场站均", summary.get("onsite_average_issue_count"), "项"),
            ],
        )
        self._add_chart_slides("各单位问题与站均问题对比", self.report.get("unit_comparison") or [], name_key="unit_name", value_keys=("issue_count", "average_issue_count"), series_names=("问题总数", "站均问题"))
        previous = self.report.get("previous_month_rectification") or {}
        self._add_chart_slides("上月各单位整改情况", previous.get("units") or [], name_key="unit_name", value_keys=("unreceived_count", "pending_count", "rectified_count"), series_names=("未签收", "未整改", "已整改"), narrative=previous.get("narrative") or "")
        for mode in self.report.get("mode_summaries") or []:
            self._add_chart_slides(f"{mode.get('label')} · 单位站均问题", mode.get("units") or [], name_key="unit_name", value_keys=("average_issue_count",), series_names=("站均问题",), narrative=mode.get("narrative") or "")
        self._add_section("CHAPTER 02", "分环节汇总")
        for section in self.report.get("category_sections") or []:
            items = section.get("items") or []
            self._add_chart_slides(f"{section.get('label')} · 分类分布", items, name_key="name", narrative=section.get("narrative") or "")
        deep = self.report.get("deep_analysis") or {}
        self._add_section("CHAPTER 03", "各单位问题分析", "AI按服务环节筛选重点问题并保留原始描述与照片。")
        for unit in (deep.get("unit_analyses") or [])[:20]:
            areas = (
                unit.get("service_areas")
                or unit.get("areas")
                or unit.get("area_analyses")
                or []
            )
            issues = self._collect_nested_issues(
                areas,
                defaults={"unit_name": unit.get("unit_name")},
            )
            self._add_issue_slides(f"{unit.get('unit_name') or '单位'} · 重点问题", issues, ai=bool(unit.get("ai_generated", True)), subtitle=unit.get("summary") or "", max_issues=8)
        self._add_section("CHAPTER 04", "问题总结与下一步建议")
        self._add_analysis_cards("问题总结", deep.get("problem_summary") or [], ai=True)
        self._add_analysis_cards("下一步建议", deep.get("next_steps") or [], ai=True)

    def _build_non_oil(self):
        summary = self.report.get("summary") or {}
        self._add_section(
            "CHAPTER 01",
            "总体情况概述",
            "按非油巡检周期汇总站点覆盖、问题数量和上期整改情况。",
        )
        self._add_narrative_slide(
            "非油检查概览",
            f"{self.report.get('period_text') or ''}\n{self.report.get('scope_text') or ''}",
            metrics=[
                ("覆盖站点", summary.get("station_count"), "座"),
                ("管理单位", summary.get("unit_count"), "个"),
                ("发现问题", summary.get("total_issue_count"), "项"),
                ("站均问题", summary.get("average_issue_count"), "项"),
            ],
        )
        previous = self.report.get("previous_month_rectification") or {}
        self._add_chart_slides(
            "上期各单位整改情况",
            previous.get("units") or [],
            name_key="unit_name",
            value_keys=(
                "total_count",
                "pending_acceptance_count",
                "pending_rectification_count",
            ),
            series_names=("全部问题", "待验收", "待整改"),
            narrative=previous.get("narrative") or "",
        )
        units = self.report.get("units") or []
        self._add_table_slides(
            "巡检范围",
            ["所属单位", "单位类型", "站点数量", "站点"],
            [
                [
                    item.get("unit_name"),
                    item.get("unit_type_label"),
                    item.get("station_count"),
                    "、".join(item.get("station_names") or []),
                ]
                for item in units
            ],
            rows_per_slide=8,
        )
        self._add_chart_slides(
            "各单位问题与站均问题",
            units,
            name_key="unit_name",
            value_keys=("issue_count", "station_count", "average_issue_count"),
            series_names=("问题总数", "站点数量", "站均问题"),
            narrative=self.report.get("issue_overview_text") or "",
        )

        deep = self.report.get("deep_analysis") or {}
        self._add_section(
            "CHAPTER 02",
            "片区问题汇总",
            "AI按单位筛选具有代表性的真实问题，保留原始描述与照片。",
        )
        for unit in (deep.get("unit_highlights") or [])[:20]:
            self._add_chart_slides(
                f"{unit.get('unit_name') or '管理单位'} · 问题分类",
                unit.get("category_distribution") or [],
                narrative=unit.get("summary") or "",
            )
            self._add_issue_slides(
                f"{unit.get('unit_name') or '管理单位'} · 重点问题",
                unit.get("highlighted_issues") or [],
                ai=bool(unit.get("ai_generated")),
                subtitle=(
                    f"涉及{unit.get('station_count') or 0}座站点，共{unit.get('issue_count') or 0}项问题，"
                    f"站均{_format_metric(unit.get('average_issue_count'))}项。"
                    f"{unit.get('summary') or ''}"
                ),
                max_issues=8,
            )

        self._add_section(
            "CHAPTER 03",
            "重点问题分析",
            "结合问题频次、覆盖站点和原始描述识别典型问题。",
        )
        for typical in (deep.get("typical_issues") or [])[:8]:
            self._add_issue_slides(
                typical.get("title") or typical.get("category_name") or "典型问题",
                typical.get("issues") or [],
                ai=bool(typical.get("ai_generated")),
                subtitle=typical.get("summary") or "",
                max_issues=8,
            )

        self._add_section(
            "CHAPTER 04",
            "具体结果分析",
            "从六类非油业务问题分布与原始检查项目两个维度分析。",
        )
        self._add_chart_slides(
            "六类非油问题分布",
            self.report.get("category_distribution") or [],
            narrative=self.report.get("category_distribution_text") or "",
        )
        category_names = [
            item.get("name") for item in (self.report.get("category_distribution") or [])
        ]
        self._add_table_slides(
            "检查项目与问题分类关联",
            ["原始检查项目", *category_names, "合计"],
            [
                [
                    item.get("source_project"),
                    *[
                        (item.get("category_counts") or {}).get(category_name, 0)
                        for category_name in category_names
                    ],
                    item.get("total_count"),
                ]
                for item in (self.report.get("project_matrix") or [])
            ],
            rows_per_slide=8,
        )
        self._add_analysis_cards(
            "分析方法",
            deep.get("analysis_method") or [],
            kicker="ANALYSIS METHOD",
            ai=False,
        )

        self._add_section(
            "CHAPTER 05",
            "检查结果分析",
            "AI结合真实问题完成归因分析并形成可执行的改善建议。",
        )
        self._add_analysis_cards(
            "归因分析框架",
            deep.get("attribution_analysis") or [],
            ai=bool(deep.get("ai_generated")),
        )
        self._add_analysis_cards(
            "问题改善建议",
            deep.get("improvement_suggestions") or [],
            ai=bool(deep.get("ai_generated")),
        )

    def _add_ending(self):
        slide = self._blank_slide(background=INK, footer=False)
        self._add_text(slide, "END", 0.84, 1.48, 2.0, 0.38, size=11, color=RGBColor(94, 234, 212), bold=True)
        self._add_text(slide, "数据驱动检查\n闭环推动改进", 0.8, 2.0, 7.8, 1.75, size=34, color=WHITE, bold=True, valign=MSO_ANCHOR.MIDDLE)
        self._add_text(slide, "本演示文稿由业务督导中心数智管理平台生成", 0.84, 5.72, 8.0, 0.4, size=12, color=RGBColor(148, 163, 184))


def build_inspection_report_presentation(report_type, report, storage_root, output_path):
    builder = InspectionReportPresentation(
        report_type=report_type,
        report=report,
        storage_root=storage_root,
    )
    return builder.build(output_path)
