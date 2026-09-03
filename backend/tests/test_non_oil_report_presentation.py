import tempfile
import unittest
from pathlib import Path
from zipfile import ZipFile

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.chart import XL_CHART_TYPE

from ai_prompts import (
    build_non_oil_category_classification_prompt,
    build_non_oil_key_issue_classification_prompt,
    build_non_oil_report_insight_prompt,
)
from non_oil_report_presentation import (
    CANVAS_SIZE,
    CATEGORY_DISPLAY_NAMES,
    TEMPLATE_FILE,
    UNIT_ORDER,
    _add_pie_chart,
    _edit_analysis_overview,
    _edit_scope_slide,
    _edit_overview_slide,
    _edit_rectification_slide,
    _add_unit_slide,
    _iter_shapes_recursive,
    _normalize_presentation_fonts,
    _remove_presentation_comments,
    build_non_oil_template_presentation,
    copy_existing_non_oil_presentation,
)


class NonOilReportPresentationTest(unittest.TestCase):
    def make_report(self):
        categories = [
            {"name": "店销商品摆放情况", "count": 3},
            {"name": "便利店卫生情况", "count": 2},
        ]
        return {
            "month": "2026-07",
            "period_text": "7月巡检，2026年7月1日-2026年7月31日",
            "scope_text": "7月非油现场抽检2座站点，涉及1个管理片区。",
            "unit_overview_text": "本次非油巡检共发现5项问题。",
            "summary": {"total_issue_count": 5, "category_count": 2},
            "previous_month_rectification": {"narrative": "上期暂无待整改问题。", "units": []},
            "category_distribution": categories,
            "key_issue_count": 2,
            "key_issue_percentage": 40,
            "key_issue_summary": {
                "selected_count": 2,
                "percentage_of_all": 40,
                "distribution": [
                    {"name": "重点商品", "count": 1, "percentage": 50},
                    {"name": "月度盘点", "count": 1, "percentage": 50},
                    {"name": "商品过期", "count": 0, "percentage": 0},
                    {"name": "团购问题", "count": 0, "percentage": 0},
                ],
                "details": [
                    {
                        "name": "重点商品",
                        "count": 1,
                        "percentage_of_all": 20,
                        "relationship": [
                            {
                                "name": "商品盘点",
                                "source_name": "商品订单、入库、盘点等情况",
                                "count": 1,
                                "percentage": 100,
                            }
                        ],
                        "issues": [
                            {
                                "issue_id": 1,
                                "source_project": "现场抽盘",
                                "station_name": "测试一站",
                                "description": "软中华盘亏10包。",
                                "issue_photo": "",
                            }
                        ],
                    },
                    {
                        "name": "月度盘点",
                        "count": 1,
                        "percentage_of_all": 20,
                        "relationship": [
                            {
                                "name": "商品盘点",
                                "source_name": "商品订单、入库、盘点等情况",
                                "count": 1,
                                "percentage": 100,
                            }
                        ],
                        "issues": [],
                    },
                ],
            },
            "category_details": [
                {
                    "source_name": "店销商品摆放情况",
                    "name": "商品摆放",
                    "count": 3,
                    "percentage": 60,
                    "features": [{"name": "价签标识", "count": 2, "example": "商品无价签"}],
                    "issues": [{"station_name": "测试一站", "description": "商品无价签", "issue_photo": ""}],
                },
                {
                    "source_name": "便利店卫生情况",
                    "name": "便利店卫生",
                    "count": 2,
                    "percentage": 40,
                    "features": [{"name": "货架积灰", "count": 2, "example": "货架有灰尘"}],
                    "issues": [{"station_name": "测试二站", "description": "货架有灰尘", "issue_photo": ""}],
                },
            ],
            "deep_analysis": {
                "core_findings": [
                    {"title": "问题集中化", "content": "问题主要集中在商品摆放。"},
                    {"title": "区域差异化", "content": "不同站点存在差异。"},
                    {"title": "风险同质化", "content": "价签问题重复出现。"},
                ],
                "attribution_analysis": [
                    {"title": "流程执行", "content": "基础流程落地不足。"},
                    {"title": "风险认知", "content": "风险认知需提升。"},
                    {"title": "监督机制", "content": "监督复核仍需强化。"},
                ],
                "action_priorities": [
                    {"title": "固化流程", "content": "固化基础检查流程。"},
                    {"title": "优化培训", "content": "开展案例化培训。"},
                    {"title": "闭环复核", "content": "建立整改复核机制。"},
                ],
            },
            "units": [
                {
                    "unit_type": "region",
                    "unit_name": "浦东",
                    "station_count": 2,
                    "station_names": ["测试一站", "测试二站"],
                    "issue_count": 5,
                    "average_issue_count": 2.5,
                    "percentage": 100,
                    "category_distribution": categories,
                    "station_issue_rows": [
                        {
                            "station_name": "测试一站",
                            "issue_count": 3,
                            "category_counts": {"店销商品摆放情况": 2, "便利店卫生情况": 1},
                        },
                        {
                            "station_name": "测试二站",
                            "issue_count": 2,
                            "category_counts": {"店销商品摆放情况": 1, "便利店卫生情况": 1},
                        },
                    ],
                }
            ],
        }

    def test_ai_prompts_remain_independent_and_complete(self):
        insight_prompt = build_non_oil_report_insight_prompt({"issues": [{"issue_id": 7}]})
        classification_prompt = build_non_oil_category_classification_prompt(
            {"allowed_categories": ["便利店卫生情况"], "issues": [{"issue_id": 7}]}
        )
        self.assertIn("typical_issues", insight_prompt)
        self.assertIn("attribution_analysis", insight_prompt)
        self.assertIn("classifications", classification_prompt)
        self.assertIn("不能输出‘其他’", classification_prompt)
        key_prompt = build_non_oil_key_issue_classification_prompt(
            {"allowed_categories": ["重点商品", "不纳入重点问题"], "issues": [{"issue_id": 7}]}
        )
        self.assertIn("不纳入重点问题", key_prompt)
        self.assertIn("不得强制归类", key_prompt)

    def make_dense_report(self):
        report = self.make_report()
        names = list(CATEGORY_DISPLAY_NAMES)
        units = []
        for index, name in enumerate(UNIT_ORDER):
            stations = [
                {
                    "station_name": f"{name}第{number + 1}站",
                    "category_counts": {
                        category: (index + number + offset) % 4
                        for offset, category in enumerate(names)
                    },
                }
                for number in range(6)
            ]
            for station in stations:
                station["issue_count"] = sum(station["category_counts"].values())
            count = sum(station["issue_count"] for station in stations)
            units.append({
                "unit_name": name,
                "unit_type": "region" if index < 8 else "holding",
                "station_count": len(stations),
                "station_names": [station["station_name"] for station in stations],
                "station_issue_rows": stations,
                "issue_count": count,
                "average_issue_count": count / len(stations),
                "category_distribution": [
                    {"name": category, "count": sum(s["category_counts"][category] for s in stations)}
                    for category in names
                ],
            })
        total = sum(unit["issue_count"] for unit in units)
        for unit in units:
            unit["percentage"] = unit["issue_count"] / total * 100
        report["units"] = units
        report["summary"].update(total_issue_count=total, station_count=108)
        report["scope_text"] = (
            f"7月非油现场与团购检查覆盖108座站点，涉及8个管理片区（{'、'.join(UNIT_ORDER[:8])}）"
            f"和10个合资公司（{'、'.join(UNIT_ORDER[8:])}）。"
        )
        report["category_distribution"] = [
            {"name": name, "count": sum(unit["category_distribution"][index]["count"] for unit in units)}
            for index, name in enumerate(names)
        ]
        report["previous_month_rectification"] = {
            "narrative": "8月各片区整改情况：累计发现问题54项，待验收18项，待整改0项。",
            "units": [
                {"unit_name": name, "total_count": 3, "pending_acceptance_count": 1, "pending_rectification_count": 0}
                for name in UNIT_ORDER
            ],
        }
        return report

    def test_column_charts_use_editable_data_tables_without_bar_labels(self):
        presentation = Presentation(TEMPLATE_FILE)
        report = self.make_dense_report()
        _edit_rectification_slide(presentation.slides[3], report)
        _edit_overview_slide(presentation.slides[8], report)
        unit_slide = _add_unit_slide(presentation, presentation.slides[9], report["units"][0])
        for slide in (presentation.slides[3], presentation.slides[8], unit_slide):
            chart_shape = next(
                shape for shape in slide.shapes
                if shape.has_chart and shape.chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED
            )
            chart = chart_shape.chart
            self.assertFalse(chart.plots[0].has_data_labels)
            self.assertFalse(chart.has_legend)
            self.assertEqual(len(chart._chartSpace.xpath(".//c:dTable")), 1)
            self.assertEqual(chart._chartSpace.xpath(".//c:dTable/c:showKeys")[0].get("val"), "1")
            self.assertTrue(chart.part.chart_workbook.xlsx_part.blob)
            self.assertLessEqual(chart_shape.top + chart_shape.height, presentation.slide_height)
        overview = next(
            shape.chart for shape in presentation.slides[8].shapes
            if shape.has_chart and shape.chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED
        )
        self.assertEqual([category.label for category in overview.plots[0].categories], UNIT_ORDER)
        self.assertEqual(list(overview.series[0].values), [unit["average_issue_count"] for unit in report["units"]])
        rectification = next(shape.chart for shape in presentation.slides[3].shapes if shape.has_chart)
        self.assertEqual(list(rectification.series[2].values), [0.0] * len(UNIT_ORDER))

    def assert_percentage_only_pie(self, chart):
        self.assertEqual(chart.chart_type, XL_CHART_TYPE.PIE)
        plot = chart.plots[0]
        self.assertTrue(plot.has_data_labels)
        labels = plot.data_labels
        self.assertTrue(labels.show_percentage)
        self.assertFalse(labels.show_value)
        self.assertFalse(labels.show_category_name)
        self.assertFalse(labels.show_series_name)
        self.assertFalse(labels.show_legend_key)
        self.assertEqual(chart._chartSpace.xpath(".//c:dLbls/c:showVal")[0].get("val"), "0")
        self.assertEqual(chart._chartSpace.xpath(".//c:dLbls/c:showPercent")[0].get("val"), "1")
        self.assertTrue(chart.has_legend)
        self.assertTrue(chart.part.chart_workbook.xlsx_part.blob)

    def test_pie_labels_show_only_percentages_without_changing_source_counts(self):
        presentation = Presentation(TEMPLATE_FILE)
        chart = _add_pie_chart(
            presentation.slides[8],
            (0, 0, presentation.slide_width // 2, presentation.slide_height // 2),
            [
                {"name": "商品摆放", "count": 3},
                {"name": "便利店卫生", "count": 2},
                {"name": "仓库管理", "count": 0},
            ],
            "非油检查问题分布",
        )
        self.assert_percentage_only_pie(chart)
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "percent-only.pptx"
            presentation.save(output)
            saved_chart = Presentation(output).slides[8].shapes[-1].chart
            self.assert_percentage_only_pie(saved_chart)
            self.assertEqual(list(saved_chart.series[0].values), [3.0, 2.0])
            self.assertEqual(
                [category.label for category in saved_chart.plots[0].categories],
                ["商品摆放", "便利店卫生"],
            )

    def test_summary_has_distinct_bullets_and_bold_labels(self):
        presentation = Presentation(TEMPLATE_FILE)
        report = self.make_dense_report()
        _edit_overview_slide(presentation.slides[8], report)
        unit_slide = _add_unit_slide(presentation, presentation.slides[9], report["units"][0])
        for slide, keyword, count in (
            (presentation.slides[8], "片区（分公司）总数", 5),
            (unit_slide, "涉及站点数", 3),
        ):
            shape = next(s for s in slide.shapes if s.has_text_frame and keyword in s.text)
            self.assertEqual(len(shape.text_frame.paragraphs), count)
            self.assertLessEqual(shape.top + shape.height, presentation.slide_height)
            for paragraph in shape.text_frame.paragraphs:
                self.assertEqual(len(paragraph._p.xpath("./a:pPr/a:buChar")), 1)
            self.assertTrue(shape.text_frame.paragraphs[0].runs[0].font.bold)
            self.assertFalse(shape.text_frame.paragraphs[0].runs[1].font.bold)
            for chart in (s for s in slide.shapes if s.has_chart):
                self.assertLess(chart.top + chart.height, shape.top)
        summary = next(s for s in presentation.slides[8].shapes if s.has_text_frame and "片区（分公司）总数" in s.text)
        self.assertIn("8个管理片区和10个合资公司", summary.text)
        self.assertEqual(str(summary.text_frame.paragraphs[3].runs[0].font.color.rgb), "C00000")
        self.assertIn("《7月非油检查问题清单》", summary.text)

    def assert_analysis_metrics(self, slide, total, selected, percentage):
        for label, expected in (
            ("问题总数", total),
            ("重点问题数", selected),
            ("重点问题占比", percentage),
            ("问题覆盖领域", "20+"),
        ):
            shape = next(
                shape for shape in _iter_shapes_recursive(slide.shapes)
                if shape.has_text_frame and label in shape.text
            )
            self.assertEqual(shape.text.splitlines()[0], str(expected))

    def test_analysis_overview_uses_actual_counts_and_recalculates_percentage(self):
        for total, selected, percentage in (
            (23, 7, "30.4%"),
            (5, 2, "40%"),
            (5, 0, "0%"),
            (0, 0, "0%"),
            (17, 17, "100%"),
            (10000, 1234, "12.3%"),
        ):
            with self.subTest(total=total, selected=selected):
                presentation = Presentation(TEMPLATE_FILE)
                slide = presentation.slides[21]
                report = self.make_report()
                report["summary"]["total_issue_count"] = total
                report["key_issue_summary"].update(
                    selected_count=selected, percentage_of_all=99,
                )
                report["key_issue_count"] = 999
                report["key_issue_percentage"] = 99
                _edit_analysis_overview(slide, report)
                self.assert_analysis_metrics(slide, total, selected, percentage)
                first_result = slide._element.xml
                _edit_analysis_overview(slide, report)
                self.assertEqual(slide._element.xml, first_result)

    def test_analysis_overview_supports_legacy_count_fields(self):
        presentation = Presentation(TEMPLATE_FILE)
        report = self.make_report()
        report.pop("key_issue_summary")
        _edit_analysis_overview(presentation.slides[21], report)
        self.assert_analysis_metrics(presentation.slides[21], 5, 2, "40%")

    def test_analysis_method_and_card_styles_are_preserved(self):
        presentation = Presentation(TEMPLATE_FILE)
        slide = presentation.slides[21]
        source_shapes = {
            shape.shape_id: shape._element.xml
            for shape in _iter_shapes_recursive(slide.shapes)
            if shape.has_text_frame
        }
        metric_labels = ("问题总数", "重点问题数", "重点问题占比")
        metric_styles = {
            shape.shape_id: [
                run._r.rPr.xml
                for paragraph in shape.text_frame.paragraphs
                for run in paragraph.runs
            ]
            for shape in _iter_shapes_recursive(slide.shapes)
            if shape.has_text_frame and any(label in shape.text for label in metric_labels)
        }
        _edit_analysis_overview(slide, self.make_report())
        for shape in _iter_shapes_recursive(slide.shapes):
            if not shape.has_text_frame:
                continue
            if shape.shape_id in metric_styles:
                self.assertEqual([
                    run._r.rPr.xml
                    for paragraph in shape.text_frame.paragraphs
                    for run in paragraph.runs
                ], metric_styles[shape.shape_id])
            else:
                self.assertEqual(shape._element.xml, source_shapes[shape.shape_id])

        method = next(s for s in slide.shapes if s.has_text_frame and "随着非油业务" in s.text)
        self.assertEqual([p.text for p in method.text_frame.paragraphs], [
            "随着非油业务在加油站利润结构中的占比提升，便利店的运营效率与合规成为管理的重点。"
            "本次分析旨在识别便利店运营流程中的潜在风险，使用的方法包括：",
            "定量评估：统计各类问题发生频次，识别高风险问题领域",
            "关联分析：建立重点问题与检查项目的映射关系，掌握重点问题的发生阶段",
            "风险识别：结合外部信息识别典型问题的风险及影响，依据风险与频次展示典型问题",
            "归因优化：探究数据之后的共性缺失，提出流程化改善建议",
        ])
        for paragraph in method.text_frame.paragraphs[1:]:
            self.assertEqual(len(paragraph._p.xpath("./a:pPr/a:buChar")), 1)
            self.assertTrue(paragraph.runs[0].font.bold)
            self.assertFalse(bool(paragraph.runs[1].font.bold))

    def test_removing_comments_is_idempotent_and_does_not_modify_template(self):
        presentation = Presentation(TEMPLATE_FILE)
        with ZipFile(TEMPLATE_FILE) as original:
            self.assertTrue(any("comments/" in name for name in original.namelist()))

        _remove_presentation_comments(presentation)
        _remove_presentation_comments(presentation)
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "clean.pptx"
            presentation.save(output)
            with ZipFile(output) as deck:
                self.assertFalse(any("comment" in name.lower() for name in deck.namelist()))
                for name in deck.namelist():
                    if name.endswith(".rels"):
                        self.assertNotIn(b"/comments", deck.read(name))
                        self.assertNotIn(b"/commentAuthors", deck.read(name))
            self.assertEqual(len(Presentation(output).slides), len(presentation.slides))
        with ZipFile(TEMPLATE_FILE) as original:
            self.assertTrue(any("comments/" in name for name in original.namelist()))

    def test_export_of_older_report_also_removes_comments(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir) / "export.pptx"
            result = copy_existing_non_oil_presentation({
                "presentation": {"ppt_path": str(TEMPLATE_FILE), "slide_count": 45},
            }, output)
            self.assertEqual(result["slide_count"], 45)
            with ZipFile(output) as package:
                self.assertFalse(any("comment" in name.lower() for name in package.namelist()))
            with ZipFile(TEMPLATE_FILE) as original:
                self.assertTrue(any("comments/" in name for name in original.namelist()))

    def test_preview_images_and_export_share_the_same_slide_set(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            slides_dir = root / "report_presentations" / "task-1"
            ppt_path = slides_dir / "report.pptx"
            result = build_non_oil_template_presentation(
                self.make_report(),
                slides_dir,
                ppt_path,
            )
            self.assertEqual(result["slide_count"], 26)
            self.assertEqual(len(list(slides_dir.glob("slide-*.jpg"))), 26)
            with Image.open(slides_dir / "slide-01.jpg") as image:
                self.assertEqual(image.size, CANVAS_SIZE)
            presentation = Presentation(ppt_path)
            pie_charts = [
                shape.chart
                for slide in presentation.slides
                for shape in slide.shapes
                if shape.has_chart and shape.chart.chart_type == XL_CHART_TYPE.PIE
            ]
            # Overview, unit pages, key-issue summary and three product pies.
            self.assertEqual(len(pie_charts), 6)
            for chart in pie_charts:
                self.assert_percentage_only_pie(chart)
            with ZipFile(ppt_path) as package:
                self.assertFalse(any("comment" in name.lower() for name in package.namelist()))
            self.assertEqual(len(presentation.slides), result["slide_count"])
            analysis_slide = next(
                slide for slide in presentation.slides
                if any(s.has_text_frame and "1. 检查总体情况" in s.text for s in slide.shapes)
            )
            self.assert_analysis_metrics(analysis_slide, 5, 2, "40%")
            normalized_template = Presentation(TEMPLATE_FILE)
            _normalize_presentation_fonts(normalized_template)
            for keyword in ("随着非油业务", "2. 分析方法"):
                expected = next(
                    s for s in normalized_template.slides[21].shapes
                    if s.has_text_frame and keyword in s.text
                )
                actual = next(s for s in analysis_slide.shapes if s.has_text_frame and keyword in s.text)
                self.assertEqual(actual._element.xml, expected._element.xml)
            self.assertTrue(
                any(shape.has_text_frame for shape in presentation.slides[0].shapes)
            )
            self.assertTrue(
                any(shape.has_table for shape in presentation.slides[4].shapes)
            )
            self.assertGreaterEqual(
                sum(
                    shape.shape_type == MSO_SHAPE_TYPE.CHART
                    for shape in presentation.slides[7].shapes
                ),
                2,
            )
            key_overview_slide = next(
                slide for slide in presentation.slides
                if any(
                    getattr(shape, "has_text_frame", False) and "3. 重点问题概述" in shape.text
                    for shape in slide.shapes
                )
            )
            self.assertTrue(any(shape.has_chart for shape in key_overview_slide.shapes))
            key_relationship_slide = next(
                slide for slide in presentation.slides
                if any(
                    getattr(shape, "has_text_frame", False) and "4. 重点问题发生环节" in shape.text
                    for shape in slide.shapes
                )
            )
            self.assertTrue(any(shape.has_table for shape in key_relationship_slide.shapes))
            self.assertGreaterEqual(
                sum(
                    shape.shape_type == MSO_SHAPE_TYPE.CHART
                    for shape in presentation.slides[8].shapes
                ),
                2,
            )
            self.assertFalse(
                len(presentation.slides[0].shapes) == 1
                and presentation.slides[0].shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE
            )

            scope_slide = presentation.slides[4]
            scope_text = next(
                shape for shape in scope_slide.shapes
                if getattr(shape, "has_text_frame", False) and "巡检期间" in shape.text
            )
            scope_table = next(shape for shape in scope_slide.shapes if shape.has_table)
            self.assertEqual(len(scope_table.table.columns), 6)
            self.assertIn("测试一站（3）", scope_table.table.cell(1, 5).text.replace("\n", ""))
            self.assertFalse(any(
                shape.has_text_frame and "检查共发现" in shape.text
                for slide in presentation.slides for shape in slide.shapes
            ))
            self.assertEqual(
                presentation.slides[5].shapes.title.text,
                Presentation(TEMPLATE_FILE).slides[6].shapes.title.text,
            )
            self.assertLessEqual(scope_text.top + scope_text.height, scope_table.top)
            self.assertLessEqual(
                scope_table.top + scope_table.height,
                presentation.slide_height,
            )

            copied_path = root / "exports" / "copy.pptx"
            copied = copy_existing_non_oil_presentation(
                {
                    "presentation": {
                        "ppt_path": "report_presentations/task-1/report.pptx",
                        "slide_count": result["slide_count"],
                    }
                },
                copied_path,
                root,
            )
            self.assertEqual(copied["slide_count"], result["slide_count"])
            self.assertEqual(copied_path.read_bytes(), ppt_path.read_bytes())

    def test_scope_slide_adapts_to_all_supported_units_without_overlap(self):
        presentation = Presentation(TEMPLATE_FILE)
        report = self.make_dense_report()
        report["scope_text"] = (
            "7月非油现场与团购检查覆盖96座站点，涉及8个管理片区"
            "（浦东、闵普徐、松金、嘉青、南汇、宝静、奉贤、崇明）和10个控（参）股单位"
            "（中油奉贤、中油同盛、中油康桥、中油农工商、中油上海、中油港汇、"
            "中石油上港、中油浦东、中油华鑫、中油中燃）。"
        )
        slide = presentation.slides[4]
        _edit_scope_slide(slide, report, presentation.slides[5])
        text_shape = next(
            shape for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and "巡检期间" in shape.text
        )
        table_shape = next(shape for shape in slide.shapes if shape.has_table)
        self.assertEqual(len(table_shape.table.rows), len(UNIT_ORDER) + 1)
        self.assertLessEqual(text_shape.top + text_shape.height, table_shape.top)
        self.assertLessEqual(
            table_shape.top + table_shape.height,
            presentation.slide_height,
        )
        self.assertEqual(sum(row.height for row in table_shape.table.rows), table_shape.height)
        self.assertLessEqual(table_shape.left + table_shape.width, presentation.slide_width)
        for index, unit in enumerate(report["units"], 1):
            row = table_shape.table.rows[index]
            self.assertEqual(row.cells[1].text.replace("\n", ""), unit["unit_name"])
            self.assertEqual(row.cells[3].text, str(unit["issue_count"]))
            self.assertEqual(row.cells[4].text, f"{unit['average_issue_count']:.1f}")
            for station in unit["station_issue_rows"]:
                self.assertIn(
                    f"{station['station_name']}（{station['issue_count']}）",
                    row.cells[5].text.replace("\n", ""),
                )
            for cell in row.cells:
                required_height = sum(
                    paragraph.line_spacing
                    for paragraph in cell.text_frame.paragraphs if paragraph.runs
                ) + cell.margin_top + cell.margin_bottom
                self.assertGreaterEqual(row.height + 2, required_height)
        self.assertGreater(len({row.height for row in list(table_shape.table.rows)[1:]}), 1)

    def test_scope_table_replaces_old_columns_and_keeps_source_style(self):
        presentation = Presentation(TEMPLATE_FILE)
        source_table = next(shape for shape in presentation.slides[5].shapes if shape.has_table)
        source_fill = source_table.table.cell(0, 0).fill.fore_color.rgb
        _edit_scope_slide(presentation.slides[4], self.make_report(), presentation.slides[5])
        tables = [shape for shape in presentation.slides[4].shapes if shape.has_table]
        self.assertEqual(len(tables), 1)
        self.assertEqual(len(tables[0].table.columns), 6)
        self.assertEqual(tables[0].table.cell(0, 0).fill.fore_color.rgb, source_fill)
        shape_ids = [shape.shape_id for shape in presentation.slides[4].shapes]
        self.assertEqual(len(shape_ids), len(set(shape_ids)))

    def test_scope_table_handles_empty_data_without_stale_template_values(self):
        presentation = Presentation(TEMPLATE_FILE)
        report = self.make_report()
        report["units"] = []
        _edit_scope_slide(presentation.slides[4], report, presentation.slides[5])
        table = next(shape.table for shape in presentation.slides[4].shapes if shape.has_table)
        self.assertEqual(len(table.rows), 2)
        self.assertEqual(table.cell(1, 1).text, "暂无数据")
        self.assertEqual(table.cell(1, 5).text, "-")


if __name__ == "__main__":
    unittest.main()
