import tempfile
import unittest
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.util import Inches

from non_oil_report_presentation import (
    TEMPLATE_FILE,
    _edit_key_issue_relationship_slide,
    _edit_monthly_inventory_slide,
    _edit_rectification_slide,
    _insert_key_issue_continuations,
    _renumber_slides,
)


class NonOilReportLayoutTest(unittest.TestCase):
    def test_relationship_table_uses_percentage_heatmap_and_blank_zeroes(self):
        prs = Presentation(TEMPLATE_FILE)
        slide = prs.slides[23]
        values = [0, 14, 29, 50, 100]
        categories = [
            ("员工服务", "员工形象及开口服务情况"),
            ("便利店卫生", "便利店卫生情况"),
            ("商品摆放", "店销商品摆放情况"),
            ("仓库管理", "仓库管理情况"),
            ("销售行为", "销售行为"),
        ]
        _edit_key_issue_relationship_slide(slide, {"key_issue_summary": {"details": [{
            "name": "重点商品", "count": 100,
            "relationship": [
                {"name": name, "source_name": source, "percentage": value}
                for (name, source), value in zip(categories, values)
            ],
        }]}})
        shape = next(s for s in slide.shapes if s.has_table)
        table = shape.table
        self.assertEqual(table.cell(0, 0).text, "")
        expected = [("", "FFFFFF", "000000"), ("14.0%", "FFFFFF", "000000"),
                    ("29.0%", "B4DBED", "000000"), ("50.0%", "00AFE8", "FFFFFF"),
                    ("100.0%", "0054B6", "FFFFFF")]
        for column, (text, fill, color) in enumerate(expected, 1):
            cell = table.cell(1, column)
            self.assertEqual(cell.text, text)
            self.assertEqual(str(cell.fill.fore_color.rgb), fill)
            if text:
                self.assertEqual(str(cell.text_frame.paragraphs[0].runs[0].font.color.rgb), color)
        for row in (0, 1, 2, 3, 4):
            self.assertEqual(str(table.cell(row, 0).fill.fore_color.rgb), "FFFFFF")
        self.assertLessEqual(shape.top + shape.height, prs.slide_height)

    def test_long_rectification_date_text_does_not_overlap_chart(self):
        prs = Presentation(TEMPLATE_FILE)
        slide = prs.slides[3]
        narrative = (
            "2026年6月15日至2026年8月14日各片区整改情况：非油现场检查与非油团购累计发现问题"
            "1242项，待验收问题1155项，待整改问题87项。请各片区尽快完成问题整改。"
        )
        _edit_rectification_slide(slide, {"previous_month_rectification": {"narrative": narrative, "units": [
            {"unit_name": "浦东", "total_count": 1242, "pending_acceptance_count": 1155, "pending_rectification_count": 87},
        ]}})
        text = next(s for s in slide.shapes if s.has_text_frame and "各片区整改情况" in s.text)
        chart = next(s for s in slide.shapes if s.has_chart)
        self.assertEqual(text.text.replace("\n", ""), narrative)
        self.assertLess(text.top + text.height, chart.top)
        self.assertLessEqual(chart.top + chart.height, prs.slide_height)

    def assert_inventory_layout(self, slide, height):
        table = next(s for s in slide.shapes if s.has_table)
        summary = next(s for s in slide.shapes if s.has_text_frame and "重点核实" in s.text)
        self.assertLess(summary.top + summary.height, table.top)
        self.assertLessEqual(table.top + table.height, height - Inches(0.27))
        self.assertEqual(table.height, sum(row.height for row in table.table.rows))
        for row in table.table.rows:
            for cell in row.cells:
                required = len(cell.text_frame.paragraphs) * cell.text_frame.paragraphs[0].line_spacing + cell.margin_top + cell.margin_bottom
                self.assertLessEqual(required, row.height)

    def test_inventory_table_remains_inside_slide_and_expands_into_continuations(self):
        prs = Presentation(TEMPLATE_FILE)
        slide = prs.slides[27]
        descriptions = [f"记录{index}：" + "库存交接记录未完整填写，盘点记录缺少复核人员签字。" * 8 for index in range(12)]
        descriptions.append("超长记录：" + "月度盘点覆盖不完整。" * 170)
        issues = [{"station_name": f"测试站{index}", "source_project": "商品订单、入库、盘点等情况", "description": text} for index, text in enumerate(descriptions)]
        continuations = _edit_monthly_inventory_slide(slide, {"key_issue_summary": {"details": [{
            "name": "月度盘点", "count": len(issues), "percentage_of_all": 25, "issues": issues,
        }]}})
        self.assertGreater(len(continuations), 1)
        _insert_key_issue_continuations(prs, continuations)
        _renumber_slides(prs)
        inventory_slides = list(prs.slides)[27:28 + len(continuations)]
        extracted = []
        for page in inventory_slides:
            self.assert_inventory_layout(page, prs.slide_height)
            table = next(s for s in page.shapes if s.has_table).table
            extracted.extend(row.cells[2].text.replace("\n", "") for row in list(table.rows)[1:])
        self.assertEqual("".join(extracted), "".join(descriptions))
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / "pagination.pptx"
            prs.save(path)
            with ZipFile(path) as package:
                self.assertEqual(len(package.namelist()), len(set(package.namelist())))
            reopened = Presentation(path)
            self.assertEqual(len(reopened.slides), 45 + len(continuations))

    def test_small_inventory_table_uses_larger_readable_type(self):
        prs = Presentation(TEMPLATE_FILE)
        slide = prs.slides[27]
        continuations = _edit_monthly_inventory_slide(slide, {"key_issue_summary": {"details": [{
            "name": "月度盘点", "count": 1, "issues": [{"station_name": "测试站", "description": "交接盘点记录缺少签字。"}],
        }]}})
        self.assertEqual(continuations, [])
        self.assert_inventory_layout(slide, prs.slide_height)
        cell = next(s for s in slide.shapes if s.has_table).table.cell(1, 2)
        self.assertEqual(cell.text_frame.paragraphs[0].runs[0].font.size.pt, 14)


if __name__ == "__main__":
    unittest.main()
